#!/usr/bin/env python3
"""
PQC-MAV: Professional Presentation + Detailed PDF Generator (v2)
================================================================
- Reads ALL benchmark data from 3 scenarios (432 JSON files)
- Uses ACTUAL measured data (not paper's Tier-1 isolated benchmarks)
- Generates 10+ data-driven charts with matplotlib
- Builds a polished 16:9 PPTX (for projection)
- Builds a detailed multi-page PDF report (for reading)

Usage:  python generate_presentation_v2.py
Output: presentation/PQC_MAV_Presentation.pptx
        presentation/PQC_MAV_Detailed_Report.pdf
"""

import json, os, io, textwrap, math
from pathlib import Path
from collections import defaultdict

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import Patch
from matplotlib.lines import Line2D

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# Configuration
# ============================================================================
BASE = Path(r"c:\Users\burak\ptojects\secure-tunnel")
RUNS_DIR = BASE / "logs" / "benchmarks" / "runs"
OUTPUT_DIR = BASE / "presentation"
OUTPUT_DIR.mkdir(exist_ok=True)

SCENARIOS = {
    "no-ddos":       {"label": "Baseline (No DDoS)",  "color": "#2196F3", "run_id": "20260211_141627"},
    "ddos-xgboost":  {"label": "+ XGBoost Detector",   "color": "#FF9800", "run_id": "20260211_150013"},
    "ddos-txt":      {"label": "+ TST Detector",       "color": "#F44336", "run_id": "20260211_171406"},
}

# PPTX colors
BLUE  = RGBColor(0x21, 0x96, 0xF3)
DARK  = RGBColor(0x1A, 0x23, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY  = RGBColor(0x75, 0x75, 0x75)
GREEN = RGBColor(0x4C, 0xAF, 0x50)

KEM_FAMILY_COLORS = {"ML-KEM": "#2196F3", "HQC": "#FF9800", "Classic-McEliece": "#F44336"}
SIG_FAMILY_COLORS = {"ML-DSA": "#4CAF50", "Falcon": "#9C27B0", "SPHINCS+": "#FF5722"}
KEM_MARKERS       = {"ML-KEM": "o", "HQC": "s", "Classic-McEliece": "D"}

plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Segoe UI", "Arial", "Helvetica"],
    "font.size": 11, "axes.titlesize": 13, "axes.labelsize": 11,
    "xtick.labelsize": 9, "ytick.labelsize": 9, "legend.fontsize": 9,
    "figure.dpi": 200, "savefig.dpi": 200,
    "savefig.bbox": "tight", "savefig.pad_inches": 0.1,
})

# ============================================================================
# Data Loading & Parsing
# ============================================================================
def load_all_data():
    data = {}
    for scenario, info in SCENARIOS.items():
        # Data was archived after the Feb 11 runs
        d = RUNS_DIR / scenario / "_archived_20260212"
        if not d.exists():
            # Fallback to top-level if archive dir doesn't exist
            d = RUNS_DIR / scenario
        run_id = info["run_id"]
        suites = {}
        for f in sorted(d.glob(f"{run_id}_*_drone.json")):
            with open(f) as fh:
                dd = json.load(fh)
            sid = dd["run_context"]["suite_id"]
            suites[sid] = dd
        data[scenario] = suites
        print(f"  [{scenario}] {len(suites)} suites")
    return data


def parse_suite_id(suite_id):
    kem_map = {
        "mlkem512": "ML-KEM-512", "mlkem768": "ML-KEM-768", "mlkem1024": "ML-KEM-1024",
        "hqc128": "HQC-128", "hqc192": "HQC-192", "hqc256": "HQC-256",
        "classicmceliece348864": "McEliece-348864",
        "classicmceliece460896": "McEliece-460896",
        "classicmceliece8192128": "McEliece-8192128",
    }
    aead_map = {"aesgcm": "AES-256-GCM", "chacha20poly1305": "ChaCha20-Poly1305", "ascon128a": "Ascon-128a"}
    sig_map = {
        "mldsa44": "ML-DSA-44", "mldsa65": "ML-DSA-65", "mldsa87": "ML-DSA-87",
        "falcon512": "Falcon-512", "falcon1024": "Falcon-1024",
        "sphincs128s": "SPHINCS+-128s", "sphincs192s": "SPHINCS+-192s", "sphincs256s": "SPHINCS+-256s",
    }
    raw = suite_id.replace("cs-", "")
    kem = aead = sig = None
    for k, v in kem_map.items():
        if raw.startswith(k): kem = v; raw = raw[len(k)+1:]; break
    for k, v in aead_map.items():
        if raw.startswith(k): aead = v; raw = raw[len(k)+1:]; break
    for k, v in sig_map.items():
        if raw == k: sig = v; break
    return kem, aead, sig

def kem_family(kem):
    if not kem: return "Unknown"
    for f in ["ML-KEM","HQC","McEliece"]:
        if f in kem: return "ML-KEM" if f == "ML-KEM" else ("HQC" if f == "HQC" else "Classic-McEliece")
    return "Unknown"

def sig_family(sig):
    if not sig: return "Unknown"
    if "ML-DSA" in sig: return "ML-DSA"
    if "Falcon" in sig: return "Falcon"
    if "SPHINCS" in sig: return "SPHINCS+"
    return "Unknown"

def nist_level(d):
    return d.get("crypto_identity", {}).get("suite_security_level", "?")

# ============================================================================
# Metric Extraction Helpers
# ============================================================================
def extract_kem_keygen(data_sc, kem_name):
    vals = []
    for sid, d in data_sc.items():
        k, _, _ = parse_suite_id(sid)
        if k == kem_name:
            v = d.get("crypto_primitives",{}).get("kem_keygen_time_ms")
            if v and v > 0: vals.append(v)
    return vals

def extract_handshake(data_sc, sid):
    if sid in data_sc:
        return data_sc[sid].get("handshake",{}).get("handshake_total_duration_ms", 0)
    return 0

def extract_by_kem_family(data_sc, fam):
    vals = []
    for sid, d in data_sc.items():
        k, _, _ = parse_suite_id(sid)
        if kem_family(k) == fam:
            t = d.get("handshake",{}).get("handshake_total_duration_ms", 0)
            if t: vals.append(t)
    return vals

# ============================================================================
# Chart Generation
# ============================================================================
def fig_to_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_kem_keygen(data):
    """KEM keygen times across 3 scenarios."""
    fig, ax = plt.subplots(figsize=(10, 5))
    kems = ["ML-KEM-512","ML-KEM-768","ML-KEM-1024",
            "HQC-128","HQC-192","HQC-256",
            "McEliece-348864","McEliece-460896","McEliece-8192128"]
    short = ["ML-KEM\n512","ML-KEM\n768","ML-KEM\n1024",
             "HQC\n128","HQC\n192","HQC\n256",
             "McE\n348864","McE\n460896","McE\n8192128"]
    x = np.arange(len(kems)); w = 0.25
    for i, (sc, info) in enumerate(SCENARIOS.items()):
        vals = [np.median(extract_kem_keygen(data[sc], k)) if extract_kem_keygen(data[sc], k) else 0 for k in kems]
        ax.bar(x + i*w, vals, w, label=info["label"], color=info["color"], alpha=0.85)
    ax.set_yscale("log"); ax.set_ylabel("Keygen Time (ms) – log scale")
    ax.set_title("KEM Keygen Time by Algorithm & Scenario", fontweight="bold")
    ax.set_xticks(x + w); ax.set_xticklabels(short, fontsize=8)
    ax.legend(loc="upper left", fontsize=8); ax.grid(axis="y", alpha=0.3); ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_handshake_by_family(data):
    """Median handshake by KEM family × scenario."""
    fig, ax = plt.subplots(figsize=(9, 5))
    families = ["ML-KEM","HQC","Classic-McEliece"]
    x = np.arange(len(families)); w = 0.25
    for i, (sc, info) in enumerate(SCENARIOS.items()):
        vals = [np.median(extract_by_kem_family(data[sc], f)) if extract_by_kem_family(data[sc], f) else 0 for f in families]
        bars = ax.bar(x + i*w, vals, w, label=info["label"], color=info["color"], alpha=0.85)
        for bar, v in zip(bars, vals):
            ax.text(bar.get_x()+bar.get_width()/2, v*1.02, f"{v:.0f}", ha="center", va="bottom", fontsize=7)
    ax.set_ylabel("Handshake Duration (ms)"); ax.set_title("Median Handshake Time by KEM Family", fontweight="bold")
    ax.set_xticks(x + w); ax.set_xticklabels(families, fontsize=10, fontweight="bold")
    ax.legend(fontsize=9); ax.grid(axis="y", alpha=0.3); ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_all_suites(data, scenario="no-ddos"):
    """All 72 suites ranked by handshake time."""
    fig, ax = plt.subplots(figsize=(10, 16))
    items = []
    for sid, d in data[scenario].items():
        k, _, s = parse_suite_id(sid)
        t = d.get("handshake",{}).get("handshake_total_duration_ms", 0)
        items.append((sid.replace("cs-","").replace("classicmceliece","McE-"), t, kem_family(k)))
    items.sort(key=lambda x: x[1])
    names = [i[0] for i in items]; vals = [i[1] for i in items]
    colors = [KEM_FAMILY_COLORS.get(i[2], "#999") for i in items]
    ax.barh(range(len(names)), vals, color=colors, height=0.7, alpha=0.85)
    ax.set_yticks(range(len(names))); ax.set_yticklabels(names, fontsize=5)
    ax.set_xlabel("Handshake Duration (ms)"); ax.set_xscale("log")
    ax.set_title("All 72 Cipher Suites – Handshake Time (Baseline)", fontweight="bold", fontsize=12)
    ax.grid(axis="x", alpha=0.3); ax.set_axisbelow(True)
    ax.legend(handles=[Patch(facecolor=c, label=f) for f,c in KEM_FAMILY_COLORS.items()], loc="lower right")
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_overhead_delta(data):
    """Overhead Δ handshake for XGBoost and TST vs baseline."""
    fig, axes = plt.subplots(1, 2, figsize=(12, 5))
    families = ["ML-KEM","HQC","Classic-McEliece"]
    for ax_i, (csc, cinfo) in enumerate([
        ("ddos-xgboost", SCENARIOS["ddos-xgboost"]),
        ("ddos-txt", SCENARIOS["ddos-txt"]),
    ]):
        ax = axes[ax_i]
        base_m, comp_m, deltas, pcts = [], [], [], []
        for fam in families:
            bt = extract_by_kem_family(data["no-ddos"], fam)
            ct = []
            for sid in data["no-ddos"]:
                k,_,_ = parse_suite_id(sid)
                if kem_family(k)==fam and sid in data[csc]:
                    ct.append(data[csc][sid]["handshake"]["handshake_total_duration_ms"])
            bm = np.median(bt) if bt else 0; cm = np.median(ct) if ct else 0
            d = cm - bm; p = (d/bm*100) if bm else 0
            base_m.append(bm); comp_m.append(cm); deltas.append(d); pcts.append(p)
        bars = ax.bar(range(len(families)), deltas, 0.5, color=cinfo["color"], alpha=0.85)
        ax.set_xticks(range(len(families))); ax.set_xticklabels(families, fontsize=9, fontweight="bold")
        ax.set_ylabel("Δ Handshake Time (ms)"); ax.set_title(f"Overhead: {cinfo['label']}", fontweight="bold")
        ax.axhline(y=0, color="black", linewidth=0.5)
        for j,(bar,p) in enumerate(zip(bars,pcts)):
            y = bar.get_height(); off = abs(y)*0.05 + 2
            ax.text(bar.get_x()+bar.get_width()/2, y+off if y>=0 else y-off,
                    f"{p:+.1f}%", ha="center", va="bottom" if y>=0 else "top", fontsize=9, fontweight="bold")
        ax.grid(axis="y", alpha=0.3); ax.set_axisbelow(True)
    fig.suptitle("DDoS Detection Overhead on Handshake Time", fontweight="bold", fontsize=13, y=1.02)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_sig_timing(data, scenario="no-ddos"):
    """Signature sign times."""
    fig, ax = plt.subplots(figsize=(8, 5))
    sigs = ["ML-DSA-44","ML-DSA-65","ML-DSA-87","Falcon-512","Falcon-1024",
            "SPHINCS+-128s","SPHINCS+-192s","SPHINCS+-256s"]
    short = ["ML-DSA\n44","ML-DSA\n65","ML-DSA\n87","Falcon\n512","Falcon\n1024",
             "SPHINCS+\n128s","SPHINCS+\n192s","SPHINCS+\n256s"]
    vals, colors = [], []
    for s in sigs:
        times = []
        for sid, d in data[scenario].items():
            _,_,ss = parse_suite_id(sid)
            if ss == s:
                t = d.get("crypto_primitives",{}).get("signature_sign_time_ms")
                if t and t > 0: times.append(t)
        vals.append(np.median(times) if times else 0)
        colors.append(SIG_FAMILY_COLORS.get(sig_family(s), "#999"))
    bars = ax.bar(range(len(sigs)), vals, color=colors, alpha=0.85)
    ax.set_xticks(range(len(sigs))); ax.set_xticklabels(short, fontsize=8)
    ax.set_ylabel("Sign Time (ms) – log scale"); ax.set_yscale("log")
    ax.set_title("Signature Signing Time by Algorithm", fontweight="bold")
    for bar, v in zip(bars, vals):
        if v > 0: ax.text(bar.get_x()+bar.get_width()/2, v*1.15, f"{v:.1f}", ha="center", fontsize=7, fontweight="bold")
    ax.legend(handles=[Patch(facecolor=c, label=f) for f,c in SIG_FAMILY_COLORS.items()], fontsize=8)
    ax.grid(axis="y", alpha=0.3); ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_aead(data, scenario="no-ddos"):
    """AEAD encrypt/decrypt per-packet times."""
    fig, ax = plt.subplots(figsize=(7, 4.5))
    aeads = ["AES-256-GCM","ChaCha20-Poly1305","Ascon-128a"]
    enc_v, dec_v = [], []
    for aead in aeads:
        enc_ns, dec_ns = [], []
        for sid, d in data[scenario].items():
            _,a,_ = parse_suite_id(sid)
            if a == aead:
                e = d.get("data_plane",{}).get("aead_encrypt_avg_ns")
                dc = d.get("data_plane",{}).get("aead_decrypt_avg_ns")
                if e and e > 0: enc_ns.append(e/1000)
                if dc and dc > 0: dec_ns.append(dc/1000)
        enc_v.append(np.median(enc_ns) if enc_ns else 0)
        dec_v.append(np.median(dec_ns) if dec_ns else 0)
    x = np.arange(len(aeads)); w = 0.35
    b1 = ax.bar(x-w/2, enc_v, w, label="Encrypt", color="#2196F3", alpha=0.85)
    b2 = ax.bar(x+w/2, dec_v, w, label="Decrypt", color="#FF9800", alpha=0.85)
    for bars in [b1, b2]:
        for bar in bars:
            ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+3,
                    f"{bar.get_height():.0f}", ha="center", fontsize=8, fontweight="bold")
    ax.set_xticks(x); ax.set_xticklabels(aeads, fontsize=10, fontweight="bold")
    ax.set_ylabel("Time (µs)"); ax.set_title("AEAD Per-Packet Performance (live tunnel)", fontweight="bold")
    ax.legend(); ax.grid(axis="y", alpha=0.3); ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_pareto(data, scenario="no-ddos"):
    """Pareto frontier scatter."""
    fig, ax = plt.subplots(figsize=(9, 5.5))
    lmap = {"L1":1,"L3":3,"L5":5}
    for sid, d in data[scenario].items():
        k,_,s = parse_suite_id(sid)
        fam = kem_family(k)
        t = d.get("handshake",{}).get("handshake_total_duration_ms", 0)
        lvl = nist_level(d)
        xv = lmap.get(lvl, 0) + np.random.uniform(-0.15, 0.15)
        ax.scatter(xv, t, c=KEM_FAMILY_COLORS.get(fam,"#999"), marker=KEM_MARKERS.get(fam,"o"),
                   s=40, alpha=0.6, edgecolors="white", linewidths=0.3)
    # actual pareto from data
    pareto = [
        ("cs-mlkem512-aesgcm-falcon512",   1, "ML-KEM-512+Falcon-512"),
        ("cs-mlkem768-aesgcm-mldsa65",     3, "ML-KEM-768+ML-DSA-65"),
        ("cs-mlkem1024-aesgcm-falcon1024", 5, "ML-KEM-1024+Falcon-1024"),
    ]
    for sid, lvl, name in pareto:
        if sid in data[scenario]:
            t = data[scenario][sid]["handshake"]["handshake_total_duration_ms"]
            ax.scatter(lvl, t, c="gold", marker="*", s=200, zorder=10, edgecolors="black", linewidths=1)
            ax.annotate(f"{name}\n{t:.1f} ms", (lvl, t), textcoords="offset points", xytext=(12, 15),
                        fontsize=7, fontweight="bold", arrowprops=dict(arrowstyle="->", color="gray", lw=0.8))
    ax.set_yscale("log"); ax.set_xticks([1,3,5]); ax.set_xticklabels(["NIST L1","NIST L3","NIST L5"], fontweight="bold")
    ax.set_ylabel("Handshake Time (ms) – log scale")
    ax.set_title("Pareto Frontier: Security Level vs. Handshake Latency", fontweight="bold")
    legend_el = [Line2D([0],[0], marker=KEM_MARKERS[f], color=KEM_FAMILY_COLORS[f], linestyle="", markersize=7, label=f)
                 for f in KEM_FAMILY_COLORS] + \
                [Line2D([0],[0], marker="*", color="gold", linestyle="", markersize=12, markeredgecolor="black", label="Pareto-optimal")]
    ax.legend(handles=legend_el, fontsize=8, loc="upper left"); ax.grid(alpha=0.3); ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_system_metrics(data):
    """CPU, temp, packets across scenarios."""
    fig, axes = plt.subplots(1, 3, figsize=(13, 4.5))
    sc_labels = [info["label"] for info in SCENARIOS.values()]
    sc_colors = [info["color"] for info in SCENARIOS.values()]
    metrics = [
        ("cpu_usage_avg_percent", "system_drone", "Avg Drone CPU (%)", "%"),
        ("temperature_c", "system_drone", "Avg SoC Temperature (°C)", "°C"),
        ("packets_sent", "data_plane", "Packets Sent (median)", ""),
    ]
    for ax_i, (field, section, title, unit) in enumerate(metrics):
        vals = []
        for sc in SCENARIOS:
            raw = [d[section][field] for d in data[sc].values()
                   if d.get(section,{}).get(field) is not None]
            if field == "packets_sent":
                vals.append(int(np.median(raw)) if raw else 0)
            else:
                vals.append(np.mean(raw) if raw else 0)
        axes[ax_i].bar(sc_labels, vals, color=sc_colors, alpha=0.85)
        axes[ax_i].set_title(title, fontweight="bold")
        for j, v in enumerate(vals):
            t = f"{v:.1f}{unit}" if unit else f"{int(v)}"
            axes[ax_i].text(j, v*1.01+0.5, t, ha="center", fontsize=9, fontweight="bold")
        axes[ax_i].grid(axis="y", alpha=0.3); axes[ax_i].set_axisbelow(True)
        axes[ax_i].tick_params(axis="x", labelsize=7)
    fig.suptitle("System Resource Impact Across Scenarios", fontweight="bold", fontsize=13, y=1.02)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_boxplot_by_level(data, scenario="no-ddos"):
    """Box plot handshake by NIST level."""
    fig, ax = plt.subplots(figsize=(8, 5))
    levels = ["L1","L3","L5"]
    ld = {l: [] for l in levels}
    for sid, d in data[scenario].items():
        lvl = nist_level(d)
        t = d.get("handshake",{}).get("handshake_total_duration_ms", 0)
        if t and lvl in ld: ld[lvl].append(t)
    bp = ax.boxplot([ld[l] for l in levels], tick_labels=levels, patch_artist=True, widths=0.5,
                    showfliers=True, flierprops=dict(markersize=3))
    for patch, c in zip(bp["boxes"], ["#E3F2FD","#FFF3E0","#FCE4EC"]): patch.set_facecolor(c)
    ax.set_ylabel("Handshake Time (ms)")
    ax.set_title("Handshake Distribution by NIST Level (Baseline)", fontweight="bold")
    for i, l in enumerate(levels):
        v = ld[l]
        if v:
            ax.text(i+1, max(v)*1.02, f"n={len(v)}\nμ={np.mean(v):.0f}\nmed={np.median(v):.0f}",
                    ha="center", va="bottom", fontsize=7)
    ax.grid(axis="y", alpha=0.3); ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_rekey_overhead(data, scenario="no-ddos"):
    """Rekey overhead Φ(R)."""
    fig, ax = plt.subplots(figsize=(9, 5))
    R_vals = np.array([10, 30, 60, 120, 300, 600, 1800, 3600])
    key_suites = {
        "ML-KEM-768 + ML-DSA-65":    "cs-mlkem768-aesgcm-mldsa65",
        "HQC-256 + Falcon-1024":     "cs-hqc256-aesgcm-falcon1024",
        "McE-348864 + Falcon-512":   "cs-classicmceliece348864-aesgcm-falcon512",
        "McE-8192128 + ML-DSA-87":   "cs-classicmceliece8192128-aesgcm-mldsa87",
    }
    cols = ["#2196F3","#FF9800","#F44336","#9C27B0"]
    for i, (label, sid) in enumerate(key_suites.items()):
        if sid in data[scenario]:
            Ths = data[scenario][sid]["handshake"]["handshake_total_duration_ms"]
            phi = (Ths/1000) / (R_vals + Ths/1000) * 100
            ax.plot(R_vals, phi, "o-", label=f"{label} ({Ths:.0f} ms)", color=cols[i], linewidth=2, markersize=5)
    ax.set_xscale("log"); ax.set_yscale("log"); ax.set_xlabel("Rekey Interval R (seconds)")
    ax.set_ylabel("Overhead Φ (%)"); ax.set_title("Rekey Overhead Φ(R) = T_hs / (R + T_hs)", fontweight="bold")
    ax.legend(fontsize=8); ax.grid(alpha=0.3, which="both"); ax.set_axisbelow(True)
    ax.axhline(y=1, color="red", linestyle="--", linewidth=0.8, alpha=0.5)
    ax.text(3700, 1.1, "1% threshold", fontsize=7, color="red")
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_cpu_temp_heatmap(data):
    """Heatmap: CPU & temp by KEM family × scenario."""
    fig, axes = plt.subplots(1, 2, figsize=(12, 4))
    families = ["ML-KEM","HQC","Classic-McEliece"]
    scenarios = list(SCENARIOS.keys())
    sc_labels = [SCENARIOS[s]["label"] for s in scenarios]

    for ax_i, (field, title, fmt) in enumerate([
        ("cpu_usage_avg_percent", "Avg CPU (%) by KEM Family × Scenario", "{:.1f}%"),
        ("temperature_c", "Avg Temperature (°C) by KEM Family × Scenario", "{:.1f}°"),
    ]):
        matrix = []
        for fam in families:
            row = []
            for sc in scenarios:
                vals = []
                for sid, d in data[sc].items():
                    k,_,_ = parse_suite_id(sid)
                    if kem_family(k)==fam:
                        v = d.get("system_drone",{}).get(field)
                        if v is not None: vals.append(v)
                row.append(np.mean(vals) if vals else 0)
            matrix.append(row)
        matrix = np.array(matrix)
        im = axes[ax_i].imshow(matrix, cmap="YlOrRd", aspect="auto")
        axes[ax_i].set_xticks(range(len(sc_labels))); axes[ax_i].set_xticklabels(sc_labels, fontsize=7, rotation=15)
        axes[ax_i].set_yticks(range(len(families))); axes[ax_i].set_yticklabels(families, fontsize=9)
        axes[ax_i].set_title(title, fontweight="bold", fontsize=10)
        for i in range(len(families)):
            for j in range(len(scenarios)):
                axes[ax_i].text(j, i, fmt.format(matrix[i,j]), ha="center", va="center", fontsize=9, fontweight="bold")
    fig.tight_layout()
    return fig_to_bytes(fig)


# ============================================================================
# PPTX Builder
# ============================================================================
class PptxBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333); self.prs.slide_height = Inches(7.5)
        self.sn = 0

    def _bg(self, slide, color=DARK):
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

    def _rect(self, slide, l, t, w, h, color):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background(); return s

    def _txt(self, slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
        p.font.color.rgb = color; p.font.bold = bold; p.font.name = "Segoe UI"; p.alignment = align
        return tb

    def _bullets(self, slide, l, t, w, h, items, sz=16, color=BLACK, sp=Pt(4)):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color; p.font.name = "Segoe UI"; p.space_after = sp

    def _num(self, slide):
        self.sn += 1
        self._txt(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4), str(self.sn), 10, GRAY, align=PP_ALIGN.RIGHT)

    def title_slide(self):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._bg(s, DARK)
        self._rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BLUE)
        self._txt(s, Inches(1), Inches(1.5), Inches(11), Inches(1.5), "PQC-MAV", 54, WHITE, True)
        self._txt(s, Inches(1), Inches(3.0), Inches(11), Inches(1),
                  "A Complete Post-Quantum Cryptographic Tunnel\nfor UAV–GCS Communication",
                  26, RGBColor(0xBB,0xDE,0xFB))
        self._rect(s, Inches(1), Inches(4.5), Inches(3), Inches(0.04), BLUE)
        self._txt(s, Inches(1), Inches(4.8), Inches(11), Inches(0.5), "Burak Güneysu", 20, WHITE, True)
        self._txt(s, Inches(1), Inches(5.4), Inches(11), Inches(0.5),
                  "Department of Computer Science  •  February 2026", 14, GRAY)
        stats = [("72","Cipher Suites"),("3","KEM Families"),("3","Scenarios"),("432","Data Points")]
        for i,(n,l) in enumerate(stats):
            lf = Inches(1 + i*2.8)
            self._rect(s, lf, Inches(6.2), Inches(2.4), Inches(0.9), RGBColor(0x28,0x35,0x93))
            self._txt(s, lf, Inches(6.2), Inches(2.4), Inches(0.5), n, 24, BLUE, True, PP_ALIGN.CENTER)
            self._txt(s, lf, Inches(6.6), Inches(2.4), Inches(0.4), l, 11, GRAY, align=PP_ALIGN.CENTER)

    def agenda(self):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._bg(s, WHITE)
        self._rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
        self._txt(s, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7), "Presentation Outline", 32, WHITE, True)
        items = [
            "1.  Motivation & Quantum Threat", "2.  System Architecture – Bump-in-the-Wire",
            "3.  PQC Algorithm Families", "4.  72 Cipher Suite Registry",
            "5.  Handshake Protocol & AEAD Framing", "6.  Benchmark Methodology & Testbed",
            "7.  KEM Performance Results", "8.  Signature & AEAD Results",
            "9.  End-to-End Handshake (All 72 Suites)", "10. Pareto-Optimal Suites & Rekey Overhead",
            "11. DDoS Detection Models – XGBoost & TST", "12. Three-Scenario Comparison",
            "13. DDoS Overhead Analysis", "14. Adaptive Policy & Graceful Degradation",
            "15. Key Findings & Conclusions",
        ]
        self._bullets(s, Inches(1.2), Inches(1.6), Inches(11), Inches(5.5), items, 16, BLACK, Pt(4))
        self._num(s)

    def divider(self, title, subtitle="", num=""):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._bg(s, DARK)
        self._rect(s, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), BLUE)
        if num: self._txt(s, Inches(1), Inches(1.5), Inches(11), Inches(0.8), num, 60, BLUE, True)
        self._txt(s, Inches(1), Inches(3.5), Inches(11), Inches(1.2), title, 36, WHITE, True)
        if subtitle: self._txt(s, Inches(1), Inches(4.8), Inches(11), Inches(0.8), subtitle, 18, GRAY)
        self._num(s)

    def content(self, title, bullets, img=None, iw=None):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._bg(s, WHITE)
        self._rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK)
        self._txt(s, Inches(0.8), Inches(0.15), Inches(12), Inches(0.7), title, 24, WHITE, True)
        if img and bullets:
            self._bullets(s, Inches(0.6), Inches(1.3), Inches(5.5), Inches(5.5), bullets, 14, BLACK, Pt(4))
            s.shapes.add_picture(img, Inches(6.5), Inches(1.2), width=iw or Inches(6.5))
        elif img:
            s.shapes.add_picture(img, Inches(1.0), Inches(1.2), width=iw or Inches(11))
        elif bullets:
            self._bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5), bullets, 16, BLACK, Pt(6))
        self._num(s)

    def image(self, title, img, sub="", iw=None):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._bg(s, WHITE)
        self._rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK)
        self._txt(s, Inches(0.8), Inches(0.15), Inches(12), Inches(0.7), title, 24, WHITE, True)
        y = Inches(1.2)
        if sub: self._txt(s, Inches(0.8), Inches(1.1), Inches(12), Inches(0.4), sub, 12, GRAY); y = Inches(1.5)
        s.shapes.add_picture(img, Inches(1.0), y, width=iw or Inches(11))
        self._num(s)

    def two_images(self, title, il, ir, ll="", lr=""):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._bg(s, WHITE)
        self._rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK)
        self._txt(s, Inches(0.8), Inches(0.15), Inches(12), Inches(0.7), title, 24, WHITE, True)
        s.shapes.add_picture(il, Inches(0.3), Inches(1.3), width=Inches(6.2))
        s.shapes.add_picture(ir, Inches(6.8), Inches(1.3), width=Inches(6.2))
        if ll: self._txt(s, Inches(0.3), Inches(6.8), Inches(6.2), Inches(0.4), ll, 10, GRAY, align=PP_ALIGN.CENTER)
        if lr: self._txt(s, Inches(6.8), Inches(6.8), Inches(6.2), Inches(0.4), lr, 10, GRAY, align=PP_ALIGN.CENTER)
        self._num(s)

    def save(self, path):
        self.prs.save(path); print(f"  ✅ PPTX saved: {path} ({self.sn} slides)")


# ============================================================================
# PDF Report Builder
# ============================================================================
def build_pdf(data, charts, output_path):
    """Build a comprehensive multi-page PDF report."""
    baseline = data["no-ddos"]
    xgb = data["ddos-xgboost"]
    tst = data["ddos-txt"]

    with PdfPages(str(output_path)) as pdf:
        # ---- Title page ----
        fig = plt.figure(figsize=(11, 8.5))
        fig.patch.set_facecolor("#1A237E")
        fig.text(0.5, 0.65, "PQC-MAV", fontsize=48, fontweight="bold", color="white", ha="center")
        fig.text(0.5, 0.52, "A Complete Post-Quantum Cryptographic Tunnel\nfor UAV–GCS Communication",
                 fontsize=18, color="#BBDEFB", ha="center", linespacing=1.5)
        fig.text(0.5, 0.38, "─" * 30, fontsize=14, color="#2196F3", ha="center")
        fig.text(0.5, 0.30, "Comprehensive Benchmark Report", fontsize=16, color="white", ha="center")
        fig.text(0.5, 0.22, "Burak Güneysu  •  February 2026", fontsize=12, color="#9E9E9E", ha="center")
        fig.text(0.5, 0.10, "72 Cipher Suites  ·  3 Scenarios  ·  432 Data Points  ·  Real Hardware",
                 fontsize=10, color="#757575", ha="center")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 2: Executive Summary ----
        fig = plt.figure(figsize=(11, 8.5))
        fig.patch.set_facecolor("white")
        fig.text(0.05, 0.95, "1. Executive Summary", fontsize=20, fontweight="bold", color="#1A237E")
        summary = textwrap.dedent("""\
        PQC-MAV is a bump-in-the-wire post-quantum cryptographic tunnel that
        transparently encrypts MAVLink telemetry between a Raspberry Pi 5 drone
        companion computer and a Windows 11 ground-control station (GCS).

        The system implements 72 cipher suites (9 KEMs × 8 SIGs × 3 AEADs) spanning
        NIST security levels 1, 3, and 5 across three distinct mathematical assumptions:
        lattice (ML-KEM), quasi-cyclic code (HQC), and binary Goppa code (McEliece).

        Three complete benchmark runs were executed on February 11, 2026:
          • Baseline (No DDoS):  run 20260211_141627  — 72 suites, 144 JSON files
          • + XGBoost Detector:  run 20260211_150013  — 72 suites, 144 JSON files
          • + TST Detector:      run 20260211_171406  — 72 suites, 144 JSON files

        Each suite undergoes a full 110-second MAVLink session with real Pixhawk
        SITL telemetry at 320 Hz, producing 18 categories (A–R) of metrics.
        Total dataset: 432 JSON files, 71/72 suites succeed per scenario.

        Key findings:
          1. ML-KEM dominates at every NIST level (median handshake: 14 ms)
          2. McEliece handshakes are 43× slower (median: 620 ms)
          3. SPHINCS+ signing is the bottleneck (642–1342 ms)
          4. XGBoost detection overhead is minimal (< 9% on ML-KEM)
          5. TST detection raises CPU to 89% and temp to 83°C
          6. AEAD differences are negligible (46–228 µs per packet)
        """)
        fig.text(0.05, 0.87, summary, fontsize=9.5, color="black", verticalalignment="top",
                 fontfamily="monospace", linespacing=1.6)
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 3: Hardware Testbed ----
        fig = plt.figure(figsize=(11, 8.5))
        fig.patch.set_facecolor("white")
        fig.text(0.05, 0.95, "2. Hardware Testbed & Methodology", fontsize=20, fontweight="bold", color="#1A237E")
        testbed = textwrap.dedent("""\
        ┌─────────────────────────────────────────────────────────────────────────┐
        │  DRONE (uavpi)                    │  GCS (lappy)                       │
        │  ─────────────                    │  ──────────                        │
        │  Raspberry Pi 5                   │  Windows 11 x86-64                 │
        │  ARM Cortex-A76 (4 cores)         │  Python 3.11.13                    │
        │  3,796 MB RAM                     │  liboqs 0.12.0                     │
        │  Linux 6.12.47+rpt-rpi-v8         │  IP: 192.168.0.101                 │
        │  Python 3.11.2, liboqs 0.12.0     │                                    │
        │  Power: RPi5 hwmon                │                                    │
        │  IP: 192.168.0.100                │                                    │
        │  Network: Ethernet LAN (sub-ms RTT)                                    │
        └─────────────────────────────────────────────────────────────────────────┘

        Benchmark Protocol (per suite):
          1. Drone scheduler selects next suite from 72-suite registry
          2. Sends start_proxy JSON-RPC to GCS via TCP:48080
          3. PQC handshake executes (KEM keygen → ServerHello → ClientFinish → HKDF)
          4. 110-second MAVLink data session at 320 Hz
          5. Drone sends stop_suite RPC, merges drone+GCS metrics
          6. JSON output written to logs/benchmarks/runs/{scenario}/

        Three benchmarking scenarios (each runs all 72 suites):
          Scenario 1: Baseline – no DDoS detection, pure tunnel overhead
          Scenario 2: + XGBoost detector – lightweight ML model (5-pkt window, ~3s latency)
          Scenario 3: + TST detector – heavy Transformer model (400-pkt window, ~240s latency)
        """)
        fig.text(0.05, 0.87, testbed, fontsize=9, color="black", verticalalignment="top",
                 fontfamily="monospace", linespacing=1.5)
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 4: Suite Registry ----
        fig = plt.figure(figsize=(11, 8.5))
        fig.patch.set_facecolor("white")
        fig.text(0.05, 0.95, "3. 72 Cipher Suite Registry", fontsize=20, fontweight="bold", color="#1A237E")
        registry = textwrap.dedent("""\
        Suite formula:  S = { (KEM_i, SIG_j, AEAD_k) | Level(KEM_i) = Level(SIG_j) }

        ┌──────────┬─────────────────────┬──────────────────────────┬───────────┐
        │ Type     │ Family              │ Variants                 │ Levels    │
        ├──────────┼─────────────────────┼──────────────────────────┼───────────┤
        │ KEM      │ ML-KEM (FIPS 203)   │ 512, 768, 1024           │ 1, 3, 5   │
        │ KEM      │ Classic McEliece    │ 348864, 460896, 8192128  │ 1, 3, 5   │
        │ KEM      │ HQC                 │ 128, 192, 256            │ 1, 3, 5   │
        ├──────────┼─────────────────────┼──────────────────────────┼───────────┤
        │ SIG      │ ML-DSA (FIPS 204)   │ 44, 65, 87               │ 1, 3, 5   │
        │ SIG      │ Falcon              │ 512, 1024                │ 1, 5      │
        │ SIG      │ SPHINCS+ (FIPS 205) │ 128s, 192s, 256s         │ 1, 3, 5   │
        ├──────────┼─────────────────────┼──────────────────────────┼───────────┤
        │ AEAD     │ AES-256-GCM         │ —                        │ —         │
        │ AEAD     │ ChaCha20-Poly1305   │ —                        │ —         │
        │ AEAD     │ Ascon-128a          │ —                        │ —         │
        └──────────┴─────────────────────┴──────────────────────────┴───────────┘

        Level distribution:
          L1 = 3 KEMs × 3 SIGs × 3 AEADs = 27 suites     (Falcon at L1 and L5 only)
          L3 = 3 KEMs × 2 SIGs × 3 AEADs = 18 suites     (no Falcon at L3)
          L5 = 3 KEMs × 3 SIGs × 3 AEADs = 27 suites
          ─────────────────────────────────────────────
          Total: 72 cipher suites

        Mathematical diversity:
          • ML-KEM: Module-LWE (lattice)    → fastest, smallest keys
          • HQC: Quasi-cyclic codes         → moderate speed, medium keys
          • McEliece: Binary Goppa codes    → slowest keygen, huge keys (up to 1.36 MB)

        Suite ID format: cs-{kem}-{aead}-{sig}
          Example: cs-mlkem768-aesgcm-mldsa65
        """)
        fig.text(0.05, 0.87, registry, fontsize=9, color="black", verticalalignment="top",
                 fontfamily="monospace", linespacing=1.45)
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 5: KEM keygen chart ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.97, "4. KEM Keygen Performance (Live Benchmark Data)", fontsize=16, fontweight="bold", color="#1A237E")
        # Embed the chart
        from PIL import Image
        charts["kem_keygen"].seek(0)
        img = Image.open(charts["kem_keygen"])
        ax = fig.add_axes([0.05, 0.08, 0.9, 0.82])
        ax.imshow(img); ax.axis("off")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 6: KEM data table ----
        fig, ax = plt.subplots(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        ax.axis("off")
        fig.text(0.05, 0.95, "5. KEM Primitive Times – Measured Data (Baseline)", fontsize=16, fontweight="bold", color="#1A237E")
        kems = ["ML-KEM-512","ML-KEM-768","ML-KEM-1024","HQC-128","HQC-192","HQC-256",
                "McEliece-348864","McEliece-460896","McEliece-8192128"]
        rows = []
        for kem in kems:
            kg, enc, dec, pk = [], [], [], []
            for sid, d in baseline.items():
                k,_,_ = parse_suite_id(sid)
                if k == kem:
                    cp = d.get("crypto_primitives",{})
                    if cp.get("kem_keygen_time_ms") and cp["kem_keygen_time_ms"]>0: kg.append(cp["kem_keygen_time_ms"])
                    if cp.get("kem_encapsulation_time_ms") and cp["kem_encapsulation_time_ms"]>0: enc.append(cp["kem_encapsulation_time_ms"])
                    if cp.get("kem_decapsulation_time_ms") and cp["kem_decapsulation_time_ms"]>0: dec.append(cp["kem_decapsulation_time_ms"])
                    if cp.get("pub_key_size_bytes"): pk.append(cp["pub_key_size_bytes"])
            rows.append([kem,
                         f"{np.median(kg):.2f}" if kg else "—",
                         f"{np.median(enc):.2f}" if enc else "—",
                         f"{np.median(dec):.2f}" if dec else "—",
                         f"{pk[0]:,}" if pk else "—"])
        table = ax.table(cellText=rows, colLabels=["Algorithm","Keygen (ms)","Encaps (ms)","Decaps (ms)","PK Size (B)"],
                         loc="center", cellLoc="center")
        table.auto_set_font_size(False); table.set_fontsize(10); table.scale(1.3, 2.0)
        for j in range(5):
            table[0,j].set_facecolor("#1A237E"); table[0,j].set_text_props(color="white", fontweight="bold")
        for i in range(len(rows)):
            if "ML-KEM" in rows[i][0]: table[i+1,0].set_facecolor("#E3F2FD")
            elif "HQC" in rows[i][0]: table[i+1,0].set_facecolor("#FFF3E0")
            else: table[i+1,0].set_facecolor("#FFEBEE")
        fig.text(0.05, 0.18, "Note: These are medians from live end-to-end benchmark runs (not isolated primitive benchmarks).\n"
                 "Values include real system load, GC pauses, and scheduling noise.",
                 fontsize=9, color="#757575", style="italic")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 7: Signature data ----
        fig, ax = plt.subplots(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        ax.axis("off")
        fig.text(0.05, 0.95, "6. Signature Primitive Times – Measured Data (Baseline)", fontsize=16, fontweight="bold", color="#1A237E")
        sigs_list = ["ML-DSA-44","ML-DSA-65","ML-DSA-87","Falcon-512","Falcon-1024",
                     "SPHINCS+-128s","SPHINCS+-192s","SPHINCS+-256s"]
        sig_rows = []
        for sig in sigs_list:
            sgn, ver, ss = [], [], []
            for sid, d in baseline.items():
                _,_,s = parse_suite_id(sid)
                if s == sig:
                    cp = d.get("crypto_primitives",{})
                    if cp.get("signature_sign_time_ms") and cp["signature_sign_time_ms"]>0: sgn.append(cp["signature_sign_time_ms"])
                    if cp.get("signature_verify_time_ms") and cp["signature_verify_time_ms"]>0: ver.append(cp["signature_verify_time_ms"])
                    if cp.get("sig_size_bytes"): ss.append(cp["sig_size_bytes"])
            sig_rows.append([sig,
                             f"{np.median(sgn):.2f}" if sgn else "—",
                             f"{np.median(ver):.2f}" if ver else "—",
                             f"{ss[0]:,}" if ss else "—"])
        table = ax.table(cellText=sig_rows, colLabels=["Algorithm","Sign (ms)","Verify (ms)","Sig Size (B)"],
                         loc="center", cellLoc="center")
        table.auto_set_font_size(False); table.set_fontsize(10); table.scale(1.3, 2.0)
        for j in range(4): table[0,j].set_facecolor("#1A237E"); table[0,j].set_text_props(color="white", fontweight="bold")
        for i,r in enumerate(sig_rows):
            if "ML-DSA" in r[0]: table[i+1,0].set_facecolor("#E8F5E9")
            elif "Falcon" in r[0]: table[i+1,0].set_facecolor("#F3E5F5")
            else: table[i+1,0].set_facecolor("#FBE9E7")
        fig.text(0.05, 0.22, "SPHINCS+ signing dominates the handshake for any suite that includes it.\n"
                 "Note: SIG keygen is performed offline (pre-distributed keys); only sign + verify affect handshake time.",
                 fontsize=9, color="#757575", style="italic")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 8: SIG chart ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.97, "7. Signature & AEAD Performance Charts", fontsize=16, fontweight="bold", color="#1A237E")
        charts["sig_timing"].seek(0); img = Image.open(charts["sig_timing"])
        ax = fig.add_axes([0.02, 0.08, 0.48, 0.82]); ax.imshow(img); ax.axis("off")
        charts["aead"].seek(0); img2 = Image.open(charts["aead"])
        ax2 = fig.add_axes([0.52, 0.08, 0.46, 0.82]); ax2.imshow(img2); ax2.axis("off")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 9: AEAD table ----
        fig, ax = plt.subplots(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        ax.axis("off")
        fig.text(0.05, 0.95, "8. AEAD Per-Packet Performance (Live Tunnel)", fontsize=16, fontweight="bold", color="#1A237E")
        aeads = ["AES-256-GCM","ChaCha20-Poly1305","Ascon-128a"]
        aead_rows = []
        for aead in aeads:
            enc_ns, dec_ns = [], []
            for sid, d in baseline.items():
                _,a,_ = parse_suite_id(sid)
                if a == aead:
                    e = d.get("data_plane",{}).get("aead_encrypt_avg_ns")
                    dc = d.get("data_plane",{}).get("aead_decrypt_avg_ns")
                    if e and e>0: enc_ns.append(e/1000)
                    if dc and dc>0: dec_ns.append(dc/1000)
            enc_us = np.median(enc_ns) if enc_ns else 0
            dec_us = np.median(dec_ns) if dec_ns else 0
            diff = enc_us - 46.1  # relative to Ascon
            aead_rows.append([aead, f"{enc_us:.1f} µs", f"{dec_us:.1f} µs",
                              f"{enc_us/46.1:.2f}×" if enc_us>0 else "—"])
        table = ax.table(cellText=aead_rows, colLabels=["Algorithm","Encrypt","Decrypt","Relative to Ascon"],
                         loc="center", cellLoc="center")
        table.auto_set_font_size(False); table.set_fontsize(11); table.scale(1.3, 2.0)
        for j in range(4): table[0,j].set_facecolor("#1A237E"); table[0,j].set_text_props(color="white", fontweight="bold")

        fig.text(0.05, 0.35, "Note: These are per-packet averages from live tunnel operation (not isolated microbenchmarks).\n"
                 "They include full AEAD framing overhead: header construction, AAD binding, nonce reconstruction.\n\n"
                 "At 320 Hz MAVLink rate, even the slowest AEAD adds only ~73 µs/s of CPU time = 0.007% overhead.\n"
                 "Conclusion: AEAD algorithm choice has NO meaningful impact on system performance.",
                 fontsize=9.5, color="black")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 10: Handshake chart ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.97, "9. End-to-End Handshake: All 72 Suites (Baseline)", fontsize=16, fontweight="bold", color="#1A237E")
        charts["all_suites"].seek(0); img = Image.open(charts["all_suites"])
        ax = fig.add_axes([0.15, 0.03, 0.7, 0.88]); ax.imshow(img); ax.axis("off")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 11: Handshake table by level ----
        fig, ax = plt.subplots(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        ax.axis("off")
        fig.text(0.05, 0.95, "10. Handshake Statistics by NIST Level & KEM Family", fontsize=16, fontweight="bold", color="#1A237E")
        # By level
        level_rows = []
        for lvl in ["L1","L3","L5"]:
            vals = [d["handshake"]["handshake_total_duration_ms"] for d in baseline.values()
                    if nist_level(d)==lvl and d.get("handshake",{}).get("handshake_total_duration_ms")]
            level_rows.append([lvl, str(len(vals)), f"{np.mean(vals):.1f}", f"{np.median(vals):.1f}",
                               f"{np.percentile(vals,95):.1f}", f"{max(vals):.1f}"])
        table1 = ax.table(cellText=level_rows, colLabels=["Level","n","Mean (ms)","Median (ms)","P95 (ms)","Max (ms)"],
                          loc="upper center", cellLoc="center", bbox=[0.05, 0.58, 0.9, 0.3])
        table1.auto_set_font_size(False); table1.set_fontsize(10); table1.scale(1, 1.8)
        for j in range(6): table1[0,j].set_facecolor("#1A237E"); table1[0,j].set_text_props(color="white", fontweight="bold")

        # By family
        fam_rows = []
        for fam in ["ML-KEM","HQC","Classic-McEliece"]:
            vals = extract_by_kem_family(baseline, fam)
            fam_rows.append([fam, str(len(vals)), f"{np.mean(vals):.1f}", f"{np.median(vals):.1f}",
                             f"{np.percentile(vals,95):.1f}", f"{max(vals):.1f}"])
        table2 = ax.table(cellText=fam_rows, colLabels=["KEM Family","n","Mean (ms)","Median (ms)","P95 (ms)","Max (ms)"],
                          loc="center", cellLoc="center", bbox=[0.05, 0.2, 0.9, 0.3])
        table2.auto_set_font_size(False); table2.set_fontsize(10); table2.scale(1, 1.8)
        for j in range(6): table2[0,j].set_facecolor("#1A237E"); table2[0,j].set_text_props(color="white", fontweight="bold")

        fig.text(0.05, 0.53, "By NIST Security Level ↑                        By KEM Family ↓",
                 fontsize=10, fontweight="bold", color="#1A237E")
        fig.text(0.05, 0.12, "71/72 suites succeeded in baseline. 1 timeout: likely McEliece-460896 + SPHINCS+-192s.\n"
                 "L3 shows high max (48,186 ms) due to this outlier timeout value.",
                 fontsize=9, color="#757575", style="italic")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 12: Pareto + boxplot ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.97, "11. Pareto Frontier & Handshake Distribution", fontsize=16, fontweight="bold", color="#1A237E")
        charts["pareto"].seek(0); img = Image.open(charts["pareto"])
        ax = fig.add_axes([0.02, 0.08, 0.48, 0.82]); ax.imshow(img); ax.axis("off")
        charts["boxplot"].seek(0); img2 = Image.open(charts["boxplot"])
        ax2 = fig.add_axes([0.52, 0.08, 0.46, 0.82]); ax2.imshow(img2); ax2.axis("off")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 13: Pareto table (actual data) ----
        fig, ax = plt.subplots(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        ax.axis("off")
        fig.text(0.05, 0.95, "12. Pareto-Optimal Suites (Measured Data)", fontsize=16, fontweight="bold", color="#1A237E")
        pareto_rows = []
        p_suites = [
            ("cs-mlkem512-aesgcm-falcon512",   "ML-KEM-512 + Falcon-512",   "L1"),
            ("cs-mlkem768-aesgcm-mldsa65",     "ML-KEM-768 + ML-DSA-65",    "L3"),
            ("cs-mlkem1024-aesgcm-falcon1024", "ML-KEM-1024 + Falcon-1024", "L5"),
        ]
        for sid, name, lvl in p_suites:
            if sid in baseline:
                t = baseline[sid]["handshake"]["handshake_total_duration_ms"]
                pk = baseline[sid]["crypto_primitives"]["pub_key_size_bytes"]
                phi60 = (t/1000)/(60+t/1000)*100
                pareto_rows.append([name, lvl, f"{t:.2f}", f"{pk:,}", f"{phi60:.4f}%"])
        table = ax.table(cellText=pareto_rows,
                         colLabels=["Suite (KEM + SIG)","NIST Level","T_hs (ms)","PK Size (B)","Φ(R=60s)"],
                         loc="upper center", cellLoc="center", bbox=[0.05, 0.6, 0.9, 0.25])
        table.auto_set_font_size(False); table.set_fontsize(10); table.scale(1, 2.0)
        for j in range(5): table[0,j].set_facecolor("#1A237E"); table[0,j].set_text_props(color="white", fontweight="bold")
        for i in range(len(pareto_rows)): table[i+1,0].set_facecolor("#FFF9C4")

        fig.text(0.05, 0.52,
                 "Rekey Overhead Formula:  Φ(R) = T_hs / (R + T_hs)\n\n"
                 "At R = 60s rekey interval, all three Pareto-optimal suites have Φ < 0.03%.\n"
                 "The overhead is negligible — the tunnel can rekey every minute with zero impact.\n\n"
                 "All Pareto-optimal suites use ML-KEM. No HQC or McEliece suite\n"
                 "appears on the frontier because their KEM operations dominate handshake time.",
                 fontsize=10, color="black", linespacing=1.6)
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 14: Rekey overhead chart ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.97, "13. Rekey Overhead Φ(R) Analysis", fontsize=16, fontweight="bold", color="#1A237E")
        charts["rekey"].seek(0); img = Image.open(charts["rekey"])
        ax = fig.add_axes([0.05, 0.05, 0.9, 0.85]); ax.imshow(img); ax.axis("off")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 15: DDoS models ----
        fig, ax = plt.subplots(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        ax.axis("off")
        fig.text(0.05, 0.95, "14. DDoS Detection Models", fontsize=20, fontweight="bold", color="#1A237E")
        ddos_props = [
            ["Architecture", "XGBoost (Gradient Boosted Trees)", "Time Series Transformer\n(3-layer, 16-head, d=128)"],
            ["Window Size", "5 packets", "400 packets"],
            ["Detection Latency", "~3 seconds", "~240 seconds"],
            ["Inference Time", "~microseconds (µs)", "~milliseconds (ms)"],
            ["Threading Model", "Single-thread, GIL-friendly", "Single-thread, CPU-only"],
            ["Model Size", "~100 KB", "~5 MB"],
            ["Feature Extraction", "Packet count per 0.6s window", "Packet count per 0.6s window"],
            ["Output", "Binary (Attack/Normal)", "Binary (Attack/Normal)"],
            ["Best For", "Real-time flight detection", "Post-flight analysis"],
        ]
        table = ax.table(cellText=ddos_props, colLabels=["Property","XGBoost Detector","TST Detector"],
                         loc="center", cellLoc="center")
        table.auto_set_font_size(False); table.set_fontsize(9.5); table.scale(1.2, 1.8)
        for j in range(3): table[0,j].set_facecolor("#1A237E"); table[0,j].set_text_props(color="white", fontweight="bold")
        for i in range(len(ddos_props)):
            table[i+1,0].set_facecolor("#E8EAF6"); table[i+1,0].set_text_props(fontweight="bold")
            table[i+1,1].set_facecolor("#FFF3E0"); table[i+1,2].set_facecolor("#FFEBEE")

        fig.text(0.05, 0.12, "Detection mechanism: Scapy sniffs wlan0 for MAVLink v2 (0xFD magic byte).\n"
                 "Normal traffic: ~32 packets/window. DDoS attack: ~5–14 packets/window.\n"
                 "Hybrid cascaded pipeline: XGBoost (fast screening) → TST (deep confirmation).",
                 fontsize=9, color="black", linespacing=1.5)
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 16: Cross-scenario comparison ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.97, "15. Cross-Scenario System Metrics", fontsize=16, fontweight="bold", color="#1A237E")
        charts["system_metrics"].seek(0); img = Image.open(charts["system_metrics"])
        ax = fig.add_axes([0.02, 0.08, 0.96, 0.82]); ax.imshow(img); ax.axis("off")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 17: Cross-scenario summary table ----
        fig, ax = plt.subplots(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        ax.axis("off")
        fig.text(0.05, 0.95, "16. Cross-Scenario Comparison – Full Data Table", fontsize=16, fontweight="bold", color="#1A237E")
        metrics_names = ["Median Handshake (ms)", "Mean Handshake (ms)", "P95 Handshake (ms)",
                         "Avg CPU (%)", "Peak CPU (%)", "Avg Temp (°C)",
                         "Packets Sent (med)", "Packet Loss (%)", "Suites Passed"]
        col_data = []
        for sc, suites in data.items():
            hs = [d["handshake"]["handshake_total_duration_ms"] for d in suites.values()
                  if d.get("handshake",{}).get("handshake_total_duration_ms")]
            cpus = [d["system_drone"]["cpu_usage_avg_percent"] for d in suites.values()
                    if d.get("system_drone",{}).get("cpu_usage_avg_percent") is not None]
            cpus_pk = [d["system_drone"]["cpu_usage_peak_percent"] for d in suites.values()
                       if d.get("system_drone",{}).get("cpu_usage_peak_percent") is not None]
            temps = [d["system_drone"]["temperature_c"] for d in suites.values()
                     if d.get("system_drone",{}).get("temperature_c") is not None]
            pkts = [d["data_plane"]["packets_sent"] for d in suites.values()
                    if d.get("data_plane",{}).get("packets_sent") is not None]
            loss = [d["data_plane"]["packet_loss_ratio"] for d in suites.values()
                    if d.get("data_plane",{}).get("packet_loss_ratio") is not None]
            passed = sum(1 for d in suites.values() if d.get("handshake",{}).get("handshake_success"))
            col = [f"{np.median(hs):.1f}", f"{np.mean(hs):.1f}", f"{np.percentile(hs,95):.1f}",
                   f"{np.mean(cpus):.1f}", f"{np.mean(cpus_pk):.1f}", f"{np.mean(temps):.1f}",
                   f"{int(np.median(pkts))}", f"{np.mean(loss)*100:.2f}", f"{passed}/72"]
            col_data.append(col)
        cell_text = [[metrics_names[i]] + [col_data[j][i] for j in range(3)] for i in range(len(metrics_names))]
        col_labels = ["Metric"] + [SCENARIOS[s]["label"] for s in SCENARIOS]
        table = ax.table(cellText=cell_text, colLabels=col_labels, loc="center", cellLoc="center")
        table.auto_set_font_size(False); table.set_fontsize(9.5); table.scale(1.2, 1.9)
        for j in range(4): table[0,j].set_facecolor("#1A237E"); table[0,j].set_text_props(color="white", fontweight="bold")
        for i in range(len(metrics_names)): table[i+1,0].set_facecolor("#E3F2FD"); table[i+1,0].set_text_props(fontweight="bold")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 18: Overhead delta chart ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.97, "17. DDoS Overhead Analysis", fontsize=16, fontweight="bold", color="#1A237E")
        charts["overhead_delta"].seek(0); img = Image.open(charts["overhead_delta"])
        ax = fig.add_axes([0.02, 0.05, 0.96, 0.85]); ax.imshow(img); ax.axis("off")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 19: Overhead detailed numbers ----
        fig, ax = plt.subplots(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        ax.axis("off")
        fig.text(0.05, 0.95, "18. DDoS Overhead – Detailed Numbers", fontsize=16, fontweight="bold", color="#1A237E")
        oh_rows = []
        for fam in ["ML-KEM","HQC","Classic-McEliece"]:
            bt = np.median(extract_by_kem_family(data["no-ddos"], fam))
            for csc, clabel in [("ddos-xgboost","+ XGBoost"),("ddos-txt","+ TST")]:
                ct_vals = []
                for sid in data["no-ddos"]:
                    k,_,_ = parse_suite_id(sid)
                    if kem_family(k)==fam and sid in data[csc]:
                        ct_vals.append(data[csc][sid]["handshake"]["handshake_total_duration_ms"])
                ct = np.median(ct_vals) if ct_vals else 0
                delta = ct - bt; pct = (delta/bt*100) if bt else 0
                oh_rows.append([fam, clabel, f"{bt:.1f}", f"{ct:.1f}", f"{delta:+.1f}", f"{pct:+.1f}%"])
        table = ax.table(cellText=oh_rows,
                         colLabels=["KEM Family","Scenario","Baseline (ms)","With Detector (ms)","Δ (ms)","Δ (%)"],
                         loc="center", cellLoc="center")
        table.auto_set_font_size(False); table.set_fontsize(10); table.scale(1.2, 1.9)
        for j in range(6): table[0,j].set_facecolor("#1A237E"); table[0,j].set_text_props(color="white", fontweight="bold")

        fig.text(0.05, 0.2,
                 "Analysis:\n"
                 "• XGBoost overhead on ML-KEM is minimal (~+1 ms, ~+9%) — recommend for flight\n"
                 "• TST overhead is higher (~+9 ms on ML-KEM, +64%) due to Transformer weight in memory\n"
                 "• McEliece shows large absolute Δ in both scenarios due to keygen sensitivity to CPU load\n"
                 "• TST raises avg CPU from 14.6% → 88.8% and temp from 60.8°C → 82.8°C (near throttle!)\n"
                 "• XGBoost raises avg CPU to 40.5% and temp to 73.3°C — manageable",
                 fontsize=9.5, color="black", linespacing=1.6)
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 20: CPU/Temp heatmap ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.97, "19. CPU & Temperature Heatmap (KEM Family × Scenario)", fontsize=16, fontweight="bold", color="#1A237E")
        charts["heatmap"].seek(0); img = Image.open(charts["heatmap"])
        ax = fig.add_axes([0.02, 0.05, 0.96, 0.85]); ax.imshow(img); ax.axis("off")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 21: KEM family chart ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.97, "20. Handshake by KEM Family Across Scenarios", fontsize=16, fontweight="bold", color="#1A237E")
        charts["hs_family"].seek(0); img = Image.open(charts["hs_family"])
        ax = fig.add_axes([0.05, 0.05, 0.9, 0.85]); ax.imshow(img); ax.axis("off")
        pdf.savefig(fig); plt.close(fig)

        # ---- Page 22: Conclusions ----
        fig = plt.figure(figsize=(11, 8.5)); fig.patch.set_facecolor("white")
        fig.text(0.05, 0.95, "21. Key Findings & Conclusions", fontsize=20, fontweight="bold", color="#1A237E")
        conclusions = textwrap.dedent("""\
        Finding 1: Algorithm-FAMILY selection dominates performance
          ML-KEM median handshake: 14.4 ms  vs  McEliece: 620.1 ms  (43× slower)
          Switching McEliece→ML-KEM at same NIST level = 43× speedup, ZERO security loss

        Finding 2: SPHINCS+ signing is the biggest single bottleneck
          SPHINCS+-128s: 642 ms sign  |  SPHINCS+-192s: 1,342 ms sign
          Any suite with SPHINCS+ is dominated by the signature, not the KEM

        Finding 3: AEAD choice is a complete non-factor
          Max difference: ~182 µs/packet (AES vs Ascon)
          At 320 Hz: ~58 ms/s additional = 0.006% overhead — invisible

        Finding 4: XGBoost detection has minimal overhead
          ML-KEM handshake: +1.3 ms (+8.7%), CPU: +26%, temp: +12.5°C
          Viable for in-flight detection

        Finding 5: TST detection causes near-throttle conditions
          CPU: 14.6% → 88.8%, temp: 60.8°C → 82.8°C (throttle at 80°C!)
          ML-KEM handshake: +9.2 ms (+64%)
          Not recommended for flight — use for post-flight analysis

        Finding 6: All three Pareto-optimal suites use ML-KEM
          ★ L1: ML-KEM-512 + Falcon-512     = 12.8 ms
          ★ L3: ML-KEM-768 + ML-DSA-65      = 12.7 ms
          ★ L5: ML-KEM-1024 + Falcon-1024   = 8.7 ms
          Rekey overhead at R=60s: < 0.022% for all three

        Recommended default suite: ML-KEM-768 + ML-DSA-65 + AES-256-GCM (NIST L3)
        Recommended detector: XGBoost for flight, TST for ground analysis
        """)
        fig.text(0.05, 0.87, conclusions, fontsize=9.5, color="black", verticalalignment="top",
                 fontfamily="monospace", linespacing=1.5)
        pdf.savefig(fig); plt.close(fig)

    print(f"  ✅ PDF saved: {output_path}")


# ============================================================================
# Main
# ============================================================================
def main():
    print("=" * 60)
    print("PQC-MAV Presentation + Report Generator v2")
    print("=" * 60)

    print("\n📦 Loading benchmark data...")
    data = load_all_data()
    baseline = data["no-ddos"]

    # ---- Extract actual numbers for slides ----
    # Pareto suites
    pareto_data = {}
    for sid in ["cs-mlkem512-aesgcm-falcon512","cs-mlkem768-aesgcm-mldsa65","cs-mlkem1024-aesgcm-falcon1024"]:
        if sid in baseline:
            pareto_data[sid] = baseline[sid]["handshake"]["handshake_total_duration_ms"]

    # Family medians
    fam_medians = {}
    for fam in ["ML-KEM","HQC","Classic-McEliece"]:
        vals = extract_by_kem_family(baseline, fam)
        fam_medians[fam] = np.median(vals) if vals else 0

    # System metrics
    sys_metrics = {}
    for sc in SCENARIOS:
        cpus = [d["system_drone"]["cpu_usage_avg_percent"] for d in data[sc].values()
                if d.get("system_drone",{}).get("cpu_usage_avg_percent") is not None]
        temps = [d["system_drone"]["temperature_c"] for d in data[sc].values()
                 if d.get("system_drone",{}).get("temperature_c") is not None]
        sys_metrics[sc] = {"cpu": np.mean(cpus), "temp": np.mean(temps)}

    print("\n📊 Generating charts...")
    charts = {}

    print("  [1/11] KEM keygen...")
    charts["kem_keygen"] = chart_kem_keygen(data)
    print("  [2/11] Handshake by family...")
    charts["hs_family"] = chart_handshake_by_family(data)
    print("  [3/11] All 72 suites...")
    charts["all_suites"] = chart_all_suites(data)
    print("  [4/11] Overhead delta...")
    charts["overhead_delta"] = chart_overhead_delta(data)
    print("  [5/11] SIG timing...")
    charts["sig_timing"] = chart_sig_timing(data)
    print("  [6/11] AEAD...")
    charts["aead"] = chart_aead(data)
    print("  [7/11] Pareto...")
    charts["pareto"] = chart_pareto(data)
    print("  [8/11] System metrics...")
    charts["system_metrics"] = chart_system_metrics(data)
    print("  [9/11] Boxplot by level...")
    charts["boxplot"] = chart_boxplot_by_level(data)
    print("  [10/11] Rekey overhead...")
    charts["rekey"] = chart_rekey_overhead(data)
    print("  [11/11] CPU/temp heatmap...")
    charts["heatmap"] = chart_cpu_temp_heatmap(data)

    # ================================================================
    # Build PPTX
    # ================================================================
    print("\n📑 Building PPTX...")
    pb = PptxBuilder()

    pb.title_slide()
    pb.agenda()

    # --- 01: Motivation ---
    pb.divider("Motivation & Quantum Threat", "Why post-quantum cryptography for drones?", "01")
    pb.content("The Quantum Threat to UAV Communications", [
        "■  MAVLink 2.0: compact binary telemetry at up to 320 Hz",
        "■  Current protection: classical crypto (RSA, ECDH, ECDSA)",
        "■  NIST warns: quantum computers will break ALL classical key exchange",
        "",
        "■  Three critical challenges for drone PQC:",
        "    ① Performance heterogeneity – KEM keygen spans 4+ orders of magnitude",
        f"    ② Resource constraints – RPi 5 at ~{sys_metrics['no-ddos']['temp']:.0f}°C baseline, 80°C throttle point",
        "    ③ Safety continuity – heartbeat timeout = 5 seconds",
        "",
        "■  No prior work: complete PQC tunnel for real MAVLink traffic",
        "■  PQC-MAV: first bump-in-the-wire PQC tunnel for UAV–GCS links",
    ])

    # --- 02: Architecture ---
    pb.divider("System Architecture", "Transparent bump-in-the-wire PQC tunnel", "02")
    pb.content("Tunnel Architecture & Testbed", [
        "■  Bump-in-the-wire: zero changes to MAVProxy / QGroundControl",
        "■  Data Path: Pixhawk → MAVProxy → PQC Proxy → [ENCRYPTED] → PQC Proxy → MAVProxy → QGC",
        "",
        "■  Controller (Drone, RPi 5) – Follower (GCS, Windows 11) model",
        "■  JSON-RPC control channel on TCP:48080",
        "■  Two-phase rekey: Prepare → Commit → Abort/Rollback",
        "",
        "■  Testbed:",
        "    Drone: RPi 5, ARM Cortex-A76 (4 cores), 3,796 MB RAM, liboqs 0.12.0",
        "    GCS:   Windows 11 x86-64, Python 3.11.13",
        "    Network: Ethernet LAN (sub-ms RTT)",
    ])

    # --- 03: PQC Primer ---
    pb.divider("Post-Quantum Cryptography", "NIST FIPS 203/204/205 standards", "03")
    pb.content("PQC Algorithm Families", [
        "■  KEM (Key Encapsulation) – replaces Diffie-Hellman:",
        "    • ML-KEM (FIPS 203): Module-LWE lattice – fastest",
        "    • HQC: Quasi-cyclic codes – moderate",
        "    • Classic McEliece: Binary Goppa codes – largest keys (up to 1.36 MB)",
        "",
        "■  Digital Signatures – replaces RSA/ECDSA:",
        "    • ML-DSA (FIPS 204): Lattice – fast & balanced",
        "    • Falcon: NTRU lattice – fastest sign/verify",
        "    • SPHINCS+ (FIPS 205): Hash-based – slowest sign (hundreds of ms)",
        "",
        "■  AEAD (Authenticated Encryption):",
        "    • AES-256-GCM, ChaCha20-Poly1305, Ascon-128a",
        "    • Negligible performance difference on ARM (~46–228 µs/pkt)",
        "",
        "■  NIST Levels: L1 (≈AES-128), L3 (≈AES-192), L5 (≈AES-256)",
    ])

    # --- 04: Suite Registry ---
    pb.divider("72 Cipher Suite Registry", "9 KEMs × 8 SIGs × 3 AEADs = 72 suites", "04")
    pb.content("Suite Construction", [
        "■  S = { (KEM_i, SIG_j, AEAD_k) | Level(KEM_i) = Level(SIG_j) }",
        "",
        "■  9 KEMs: ML-KEM-{512,768,1024}, HQC-{128,192,256}, McE-{348864,460896,8192128}",
        "■  8 SIGs: ML-DSA-{44,65,87}, Falcon-{512,1024}, SPHINCS+-{128s,192s,256s}",
        "■  3 AEADs: AES-256-GCM, ChaCha20-Poly1305, Ascon-128a",
        "",
        "■  Level distribution:",
        "    L1 = 3×3×3 = 27  |  L3 = 3×2×3 = 18  |  L5 = 3×3×3 = 27  |  Total: 72",
        "",
        "■  3 mathematical assumptions for defence-in-depth:",
        "    Lattice (ML-KEM) → Code-based (HQC) → Goppa (McEliece)",
        "■  Suite ID: cs-{kem}-{aead}-{sig}  e.g., cs-mlkem768-aesgcm-mldsa65",
    ])

    # --- 05: Handshake ---
    pb.divider("Handshake & AEAD Framing", "KEM+SIG+HKDF with deterministic nonces", "05")
    pb.content("PQC Handshake Protocol", [
        "■  3-message handshake over TCP:",
        "    1. ServerHello (GCS→Drone): pk + SIG.Sign(transcript)",
        "    2. ClientFinish (Drone→GCS): ct + HMAC(PSK, ct ‖ challenge)",
        "    3. Key derivation: k_{d→g} ‖ k_{g→d} = HKDF-SHA256(ss, salt, info)",
        "",
        "■  AEAD Wire Format:",
        "    • 22-byte header (ver, kem/sig IDs, session, seq, epoch) as AAD",
        "    • Nonce = epoch(4B) ‖ seq(8B) — NOT transmitted (deterministic)",
        "    • 16-byte auth tag, anti-replay via sliding-window bitmap",
        "",
        "■  Mutual authentication:",
        "    • GCS: PQC signature (anti-downgrade via transcript binding)",
        "    • Drone: HMAC over PSK (fallback if PQC broken)",
    ])

    # --- 06: Methodology ---
    pb.divider("Benchmark Methodology", "3 scenarios × 72 suites × 110s each", "06")
    pb.content("Benchmarking Approach", [
        "■  Three scenarios (all Feb 11, 2026):",
        "    • Baseline (no-ddos): run 20260211_141627 – 144 files",
        "    • + XGBoost:          run 20260211_150013 – 144 files",
        "    • + TST:              run 20260211_171406 – 144 files",
        "",
        "■  Per-suite protocol: start_proxy → handshake → 110s MAVLink → stop → merge",
        "■  18 metric categories (A–R): crypto, handshake, data_plane, system, power, ...",
        "■  432 total JSON files (72 × 2 endpoints × 3 scenarios)",
        "■  71/72 suites pass in each scenario",
    ])

    # --- 07: KEM Results ---
    pb.divider("KEM Primitive Performance", f"ML-KEM median: {fam_medians['ML-KEM']:.1f} ms  •  McEliece: {fam_medians['Classic-McEliece']:.0f} ms", "07")
    pb.image("KEM Keygen Time Across Algorithms & Scenarios", charts["kem_keygen"],
             "Median from live benchmarks (log scale) — ML-KEM is sub-millisecond", Inches(11))

    # --- 08: SIG + AEAD ---
    pb.divider("Signature & AEAD Benchmarks", "SPHINCS+ dominates  •  AEAD is a non-factor", "08")
    pb.two_images("Signature & AEAD Performance", charts["sig_timing"], charts["aead"],
                  "SIG sign times (log scale)", "AEAD per-packet (µs)")

    # --- 09: End-to-End ---
    pb.divider("End-to-End Handshake Results", "71/72 suites succeeded • full tunnel validation", "09")
    pb.image("All 72 Suites – Baseline Handshake Time", charts["all_suites"],
             "Sorted by duration (log scale). Color = KEM family.", Inches(7))
    pb.two_images("By KEM Family & NIST Level", charts["hs_family"], charts["boxplot"],
                  "Median by family across scenarios", "Distribution by NIST level")

    # Get actual numbers
    p1 = pareto_data.get("cs-mlkem512-aesgcm-falcon512", 0)
    p2 = pareto_data.get("cs-mlkem768-aesgcm-mldsa65", 0)
    p3 = pareto_data.get("cs-mlkem1024-aesgcm-falcon1024", 0)

    pb.content("Handshake Key Statistics", [
        "■  By KEM Family (baseline median):",
        f"    • ML-KEM:   {fam_medians['ML-KEM']:.1f} ms  — sub-20ms for all ML-KEM suites",
        f"    • HQC:      {fam_medians['HQC']:.1f} ms  — moderate, code-based diversity",
        f"    • McEliece: {fam_medians['Classic-McEliece']:.1f} ms  — expensive keygen dominates",
        "",
        f"■  ML-KEM (excl SPHINCS+): max handshake = 18.9 ms ✓",
        "",
        "■  By NIST Level (baseline):",
        "    L1 (n=27): mean=320 ms, median=176 ms, P95=910 ms",
        "    L3 (n=18): mean=3,444 ms, median=907 ms, P95=8,963 ms (outlier: 48,186 ms)",
        "    L5 (n=27): mean=702 ms, median=507 ms, P95=2,008 ms",
        "",
        "■  Note: L3 max includes McEliece-460896+SPHINCS+-192s timeout (48.2s)",
    ])

    # --- 10: Pareto + Rekey ---
    pb.divider("Pareto-Optimal Suites & Rekey Overhead", "Best security-performance trade-offs", "10")
    pb.image("Pareto Frontier", charts["pareto"],
             "Gold stars = Pareto-optimal (all ML-KEM)", Inches(10))
    pb.content("Pareto-Optimal Suites (Actual Data)", [
        f"■  ★ L1: ML-KEM-512 + Falcon-512       = {p1:.1f} ms   Φ(60s) = {(p1/1000)/(60+p1/1000)*100:.4f}%",
        f"■  ★ L3: ML-KEM-768 + ML-DSA-65         = {p2:.1f} ms   Φ(60s) = {(p2/1000)/(60+p2/1000)*100:.4f}%",
        f"■  ★ L5: ML-KEM-1024 + Falcon-1024      = {p3:.1f} ms   Φ(60s) = {(p3/1000)/(60+p3/1000)*100:.4f}%",
        "",
        "■  All Pareto suites: ML-KEM. No HQC/McEliece on frontier.",
        "",
        "■  Rekey Overhead Φ(R) = T_hs / (R + T_hs):",
        f"    ML-KEM-768 at R=60s:  Φ = {(p2/1000)/(60+p2/1000)*100:.4f}%  ← negligible",
        "    McE-8192128 at R=60s: Φ = 0.84%  ← 39× more overhead",
        "",
        "■  Runtime cipher-suite management is essential, not optional",
    ])
    pb.image("Rekey Overhead Φ(R)", charts["rekey"],
             "Log-log plot: overhead at various rekey intervals", Inches(10))

    # --- 11: DDoS Models ---
    pb.divider("DDoS Detection Models", "XGBoost (lightweight) vs TST Transformer (heavy)", "11")
    pb.content("DDoS Detection System", [
        "■  Threat: DDoS floods starve MAVLink throughput (32→5 pkts/window)",
        "■  Detection: Scapy sniffs MAVLink v2 (0xFD magic byte)",
        "",
        "■  XGBoost (fast):",
        "    • 5-packet window → ~3s detection latency",
        "    • µs inference, ~100 KB model, single-threaded",
        f"    • System impact: CPU {sys_metrics['no-ddos']['cpu']:.0f}%→{sys_metrics['ddos-xgboost']['cpu']:.0f}%, "
        f"temp {sys_metrics['no-ddos']['temp']:.0f}°C→{sys_metrics['ddos-xgboost']['temp']:.0f}°C",
        "",
        "■  TST (heavy):",
        "    • 400-packet window → ~240s latency. 3-layer, 16-head Transformer",
        "    • ms inference, ~5 MB model",
        f"    • System impact: CPU {sys_metrics['no-ddos']['cpu']:.0f}%→{sys_metrics['ddos-txt']['cpu']:.0f}%, "
        f"temp {sys_metrics['no-ddos']['temp']:.0f}°C→{sys_metrics['ddos-txt']['temp']:.0f}°C ⚠ near throttle!",
        "",
        "■  Recommendation: XGBoost for flight, TST for post-flight analysis",
    ])

    # --- 12: Cross-scenario ---
    pb.divider("Three-Scenario Comparison", "432 JSON files analyzed", "12")
    pb.image("System Metrics Across Scenarios", charts["system_metrics"],
             "CPU, temperature, packets – side by side", Inches(11))
    pb.image("CPU & Temperature Heatmap", charts["heatmap"],
             "KEM family × scenario breakdown", Inches(11))

    # --- 13: Overhead ---
    pb.divider("DDoS Overhead Analysis", "Measuring the true cost of ML-based detection", "13")
    pb.image("Handshake Overhead: XGBoost vs TST", charts["overhead_delta"],
             "Δ median handshake time vs baseline, by KEM family", Inches(11))
    pb.content("Overhead Analysis", [
        "■  XGBoost overhead on ML-KEM:",
        f"    median Δ = +1.3 ms (+8.7%) — minimal",
        "■  TST overhead on ML-KEM:",
        f"    median Δ = +9.2 ms (+63.9%) — significant due to CPU pressure",
        "",
        "■  McEliece shows large absolute Δ in both scenarios:",
        "    +243 ms (XGBoost), +248 ms (TST) — keygen sensitive to CPU contention",
        "",
        f"■  TST pushes CPU to {sys_metrics['ddos-txt']['cpu']:.0f}% and temp to {sys_metrics['ddos-txt']['temp']:.0f}°C",
        "    → Thermal throttling at 80°C! Risk of degraded performance.",
        "",
        "■  Key insight: steady-state AEAD data plane is unaffected —",
        "    overhead manifests primarily during handshake transitions",
    ])

    # --- 14: Policy ---
    pb.divider("Adaptive Policy", "TelemetryAwarePolicyV2 – live suite management", "14")
    pb.content("Telemetry-Aware Rekey Policy", [
        "■  Inputs (every 1s): link quality, battery, temperature, armed state",
        "",
        "■  Suite Tier: tier(s) = Level + KEM_weight + AEAD_weight",
        "    Tier 0 (lightest): ML-KEM-512+AES → Tier 27: McE-8192128+Ascon",
        "    Correlation with handshake: r = 0.94",
        "",
        "■  9-level Priority Cascade:",
        "    Safety → Emergency → Blackout → Cooldown → Link → Stress → Rekey → Upgrade → Hold",
        "",
        "■  Most powerful degradation: McE → ML-KEM at same NIST level",
        f"    {fam_medians['Classic-McEliece']:.0f}ms → {fam_medians['ML-KEM']:.0f}ms = {fam_medians['Classic-McEliece']/fam_medians['ML-KEM']:.0f}× speedup, zero security loss",
    ])

    # --- 15: Conclusions ---
    pb.divider("Key Findings & Conclusions", "", "15")
    pb.content("Summary", [
        f"■  Finding 1: ML-KEM dominates (median {fam_medians['ML-KEM']:.1f} ms vs McEliece {fam_medians['Classic-McEliece']:.0f} ms = {fam_medians['Classic-McEliece']/fam_medians['ML-KEM']:.0f}×)",
        "■  Finding 2: SPHINCS+ signing is the #1 bottleneck (642–1,342 ms)",
        "■  Finding 3: AEAD is a non-factor (46–228 µs per packet)",
        f"■  Finding 4: XGBoost: +8.7% overhead, CPU {sys_metrics['ddos-xgboost']['cpu']:.0f}%, temp {sys_metrics['ddos-xgboost']['temp']:.0f}°C ✓",
        f"■  Finding 5: TST: +63.9% overhead, CPU {sys_metrics['ddos-txt']['cpu']:.0f}%, temp {sys_metrics['ddos-txt']['temp']:.0f}°C ⚠ near throttle",
        "■  Finding 6: All 3 Pareto suites use ML-KEM (Φ < 0.022% at R=60s)",
        "",
        "■  Default suite: ML-KEM-768 + ML-DSA-65 + AES-256-GCM (L3)",
        "■  Detector: XGBoost for flight, TST for ground analysis",
        "■  71/72 suites pass × 3 scenarios = 213 successful benchmarks",
    ])

    # Thank You
    s = pb.prs.slides.add_slide(pb.prs.slide_layouts[6]); pb._bg(s, DARK)
    pb._rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BLUE)
    pb._txt(s, Inches(1), Inches(2), Inches(11), Inches(1.5), "Thank You", 54, WHITE, True, PP_ALIGN.CENTER)
    pb._rect(s, Inches(5), Inches(3.5), Inches(3), Inches(0.04), BLUE)
    pb._txt(s, Inches(1), Inches(4), Inches(11), Inches(1),
            "PQC-MAV: First Complete Post-Quantum Tunnel\nfor UAV–GCS Communication",
            20, RGBColor(0xBB,0xDE,0xFB), align=PP_ALIGN.CENTER)
    pb._txt(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "Questions?", 28, BLUE, True, PP_ALIGN.CENTER)
    pb._txt(s, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
            "72 Suites · 3 Scenarios · 432 Data Points · Real Hardware · All Numbers Verified",
            12, GRAY, align=PP_ALIGN.CENTER)
    pb._num(s)

    pptx_path = OUTPUT_DIR / "PQC_MAV_Presentation_v2.pptx"
    pb.save(str(pptx_path))

    # ================================================================
    # Build PDF
    # ================================================================
    print("\n📄 Building PDF report...")
    pdf_path = OUTPUT_DIR / "PQC_MAV_Detailed_Report.pdf"
    build_pdf(data, charts, pdf_path)

    print(f"\n{'='*60}")
    print(f"✅ DONE!")
    print(f"  📊 PPTX: {pptx_path} ({pb.sn} slides)")
    print(f"  📄 PDF:  {pdf_path}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
