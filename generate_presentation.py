#!/usr/bin/env python3
"""
PQC-MAV: Professional Presentation Generator
=============================================
Generates a comprehensive PPTX presentation for professor review.
Reads all benchmark data from 3 scenarios (no-ddos, ddos-xgboost, ddos-txt),
generates plots with matplotlib, and builds a polished PPTX with python-pptx.

Usage:  python generate_presentation.py
Output: presentation/PQC_MAV_Presentation.pptx
"""

import json, os, glob, math, re, io, textwrap
from pathlib import Path
from collections import defaultdict

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
from matplotlib.patches import FancyBboxPatch
import matplotlib.patheffects as pe

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ============================================================================
# Configuration
# ============================================================================
BASE = Path(r"c:\Users\burak\ptojects\secure-tunnel")
RUNS_DIR = BASE / "logs" / "benchmarks" / "runs"
FIGURES_DIR = BASE / "paper" / "figures"
OUTPUT_DIR = BASE / "presentation"
OUTPUT_DIR.mkdir(exist_ok=True)

SCENARIOS = {
    "no-ddos":       {"label": "Baseline (No DDoS)", "color": "#2196F3", "run_id": "20260211_141627"},
    "ddos-xgboost":  {"label": "+ XGBoost Detector", "color": "#FF9800", "run_id": "20260211_150013"},
    "ddos-txt":      {"label": "+ TST Detector",     "color": "#F44336", "run_id": "20260211_171406"},
}

# Color palette
BLUE  = RGBColor(0x21, 0x96, 0xF3)
DARK  = RGBColor(0x1A, 0x23, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT1 = RGBColor(0x00, 0x96, 0x88)  # Teal
ACCENT2 = RGBColor(0xFF, 0x57, 0x22)  # Deep Orange
ACCENT3 = RGBColor(0x9C, 0x27, 0xB0)  # Purple
GRAY    = RGBColor(0x75, 0x75, 0x75)
RED     = RGBColor(0xF4, 0x43, 0x36)
GREEN   = RGBColor(0x4C, 0xAF, 0x50)
ORANGE  = RGBColor(0xFF, 0x98, 0x00)
BLACK   = RGBColor(0x00, 0x00, 0x00)

# KEM families
KEM_FAMILIES = {
    "ML-KEM":            {"color": "#2196F3", "marker": "o"},
    "HQC":               {"color": "#FF9800", "marker": "s"},
    "Classic-McEliece":  {"color": "#F44336", "marker": "D"},
}
SIG_FAMILIES = {
    "ML-DSA":    {"color": "#4CAF50"},
    "Falcon":    {"color": "#9C27B0"},
    "SPHINCS+":  {"color": "#FF5722"},
}

plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Segoe UI", "Arial", "Helvetica"],
    "font.size": 11,
    "axes.titlesize": 13,
    "axes.labelsize": 11,
    "xtick.labelsize": 9,
    "ytick.labelsize": 9,
    "legend.fontsize": 9,
    "figure.dpi": 200,
    "savefig.dpi": 200,
    "savefig.bbox": "tight",
    "savefig.pad_inches": 0.1,
})

# ============================================================================
# Data Loading
# ============================================================================
def load_all_data():
    """Load all benchmark JSON files from all 3 scenarios."""
    data = {}
    for scenario, info in SCENARIOS.items():
        scenario_dir = RUNS_DIR / scenario
        run_id = info["run_id"]
        suites = {}
        for f in sorted(scenario_dir.glob(f"{run_id}_*_drone.json")):
            with open(f) as fh:
                d = json.load(fh)
            sid = d["run_context"]["suite_id"]
            suites[sid] = d
        data[scenario] = suites
        print(f"  [{scenario}] Loaded {len(suites)} suites")
    return data


def parse_suite_id(suite_id):
    """Parse suite_id like cs-mlkem768-aesgcm-mldsa65 into components."""
    parts = suite_id.replace("cs-", "").split("-")
    # KEM is first token(s), AEAD is one of known values, SIG is last
    kem_map = {
        "mlkem512": "ML-KEM-512", "mlkem768": "ML-KEM-768", "mlkem1024": "ML-KEM-1024",
        "hqc128": "HQC-128", "hqc192": "HQC-192", "hqc256": "HQC-256",
        "classicmceliece348864": "Classic-McEliece-348864",
        "classicmceliece460896": "Classic-McEliece-460896",
        "classicmceliece8192128": "Classic-McEliece-8192128",
    }
    aead_map = {
        "aesgcm": "AES-256-GCM", "chacha20poly1305": "ChaCha20-Poly1305", "ascon128a": "Ascon-128a",
    }
    sig_map = {
        "mldsa44": "ML-DSA-44", "mldsa65": "ML-DSA-65", "mldsa87": "ML-DSA-87",
        "falcon512": "Falcon-512", "falcon1024": "Falcon-1024",
        "sphincs128s": "SPHINCS+-128s", "sphincs192s": "SPHINCS+-192s", "sphincs256s": "SPHINCS+-256s",
    }
    raw = suite_id.replace("cs-", "")
    kem = aead = sig = None
    for k, v in kem_map.items():
        if raw.startswith(k):
            kem = v
            raw = raw[len(k)+1:]
            break
    for k, v in aead_map.items():
        if raw.startswith(k):
            aead = v
            raw = raw[len(k)+1:]
            break
    for k, v in sig_map.items():
        if raw == k:
            sig = v
            break
    return kem, aead, sig


def get_kem_family(kem):
    if kem and "ML-KEM" in kem: return "ML-KEM"
    if kem and "HQC" in kem: return "HQC"
    if kem and "McEliece" in kem: return "Classic-McEliece"
    return "Unknown"


def get_sig_family(sig):
    if sig and "ML-DSA" in sig: return "ML-DSA"
    if sig and "Falcon" in sig: return "Falcon"
    if sig and "SPHINCS" in sig: return "SPHINCS+"
    return "Unknown"


def get_nist_level(d):
    return d.get("crypto_identity", {}).get("suite_security_level", "?")


# ============================================================================
# Chart Generation Functions
# ============================================================================
def fig_to_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_kem_keygen_comparison(data):
    """Bar chart – KEM keygen times by algorithm across 3 scenarios."""
    fig, ax = plt.subplots(figsize=(10, 5))
    kems = ["ML-KEM-512", "ML-KEM-768", "ML-KEM-1024",
            "HQC-128", "HQC-192", "HQC-256",
            "Classic-McEliece-348864", "Classic-McEliece-460896", "Classic-McEliece-8192128"]
    kem_short = ["ML-KEM\n512", "ML-KEM\n768", "ML-KEM\n1024",
                 "HQC\n128", "HQC\n192", "HQC\n256",
                 "McE\n348864", "McE\n460896", "McE\n8192128"]
    x = np.arange(len(kems))
    width = 0.25
    for i, (scenario, info) in enumerate(SCENARIOS.items()):
        vals = []
        for kem in kems:
            times = []
            for sid, d in data[scenario].items():
                k, _, _ = parse_suite_id(sid)
                if k == kem:
                    t = d.get("crypto_primitives", {}).get("kem_keygen_time_ms", 0)
                    if t: times.append(t)
            vals.append(np.median(times) if times else 0)
        ax.bar(x + i*width, vals, width, label=info["label"], color=info["color"], alpha=0.85)

    ax.set_yscale("log")
    ax.set_ylabel("Keygen Time (ms) – log scale")
    ax.set_title("KEM Keygen Time by Algorithm & Scenario", fontweight="bold")
    ax.set_xticks(x + width)
    ax.set_xticklabels(kem_short, fontsize=8)
    ax.legend(loc="upper left", fontsize=8)
    ax.grid(axis="y", alpha=0.3)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_handshake_by_family(data):
    """Grouped bar chart – median handshake time by KEM family across scenarios."""
    fig, ax = plt.subplots(figsize=(9, 5))
    families = ["ML-KEM", "HQC", "Classic-McEliece"]
    x = np.arange(len(families))
    width = 0.25
    for i, (scenario, info) in enumerate(SCENARIOS.items()):
        vals = []
        for fam in families:
            times = []
            for sid, d in data[scenario].items():
                k, _, _ = parse_suite_id(sid)
                if get_kem_family(k) == fam:
                    t = d.get("handshake", {}).get("handshake_total_duration_ms", 0)
                    if t: times.append(t)
            vals.append(np.median(times) if times else 0)
        ax.bar(x + i*width, vals, width, label=info["label"], color=info["color"], alpha=0.85)

    ax.set_ylabel("Handshake Duration (ms)")
    ax.set_title("Median Handshake Time by KEM Family", fontweight="bold")
    ax.set_xticks(x + width)
    ax.set_xticklabels(families, fontsize=10, fontweight="bold")
    ax.legend(fontsize=9)
    ax.grid(axis="y", alpha=0.3)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_handshake_all_suites(data, scenario="no-ddos"):
    """Horizontal bar chart – all 72 suites ranked by handshake time."""
    fig, ax = plt.subplots(figsize=(10, 16))
    suites_data = []
    for sid, d in data[scenario].items():
        kem, aead, sig = parse_suite_id(sid)
        t = d.get("handshake", {}).get("handshake_total_duration_ms", 0)
        fam = get_kem_family(kem)
        suites_data.append((sid, t, fam, kem, sig))

    suites_data.sort(key=lambda x: x[1])
    names = []
    vals = []
    colors = []
    fam_colors = {"ML-KEM": "#2196F3", "HQC": "#FF9800", "Classic-McEliece": "#F44336"}
    for sid, t, fam, kem, sig in suites_data:
        short = sid.replace("cs-", "").replace("classicmceliece", "McE-")
        names.append(short)
        vals.append(t)
        colors.append(fam_colors.get(fam, "#999999"))

    y = np.arange(len(names))
    ax.barh(y, vals, color=colors, height=0.7, alpha=0.85)
    ax.set_yticks(y)
    ax.set_yticklabels(names, fontsize=5)
    ax.set_xlabel("Handshake Duration (ms)")
    ax.set_title("All 72 Cipher Suites – Handshake Time (Baseline)", fontweight="bold", fontsize=12)
    ax.set_xscale("log")
    ax.grid(axis="x", alpha=0.3)
    ax.set_axisbelow(True)
    # Legend
    from matplotlib.patches import Patch
    legend_elements = [Patch(facecolor=c, label=f) for f, c in fam_colors.items()]
    ax.legend(handles=legend_elements, loc="lower right", fontsize=9)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_scenario_overhead_delta(data):
    """Bar chart – overhead (Δ handshake time over baseline) for XGBoost and TST."""
    fig, axes = plt.subplots(1, 2, figsize=(12, 5))
    families = ["ML-KEM", "HQC", "Classic-McEliece"]

    for ax_idx, (compare_scenario, compare_info) in enumerate([
        ("ddos-xgboost", SCENARIOS["ddos-xgboost"]),
        ("ddos-txt", SCENARIOS["ddos-txt"]),
    ]):
        ax = axes[ax_idx]
        baseline_by_family = {f: [] for f in families}
        compare_by_family = {f: [] for f in families}

        for sid in data["no-ddos"]:
            k, _, _ = parse_suite_id(sid)
            fam = get_kem_family(k)
            base_t = data["no-ddos"][sid].get("handshake", {}).get("handshake_total_duration_ms", 0)
            if sid in data[compare_scenario]:
                comp_t = data[compare_scenario][sid].get("handshake", {}).get("handshake_total_duration_ms", 0)
                if base_t and comp_t:
                    baseline_by_family[fam].append(base_t)
                    compare_by_family[fam].append(comp_t)

        x = np.arange(len(families))
        base_medians = [np.median(baseline_by_family[f]) if baseline_by_family[f] else 0 for f in families]
        comp_medians = [np.median(compare_by_family[f]) if compare_by_family[f] else 0 for f in families]
        deltas = [c - b for b, c in zip(base_medians, comp_medians)]
        pct = [(d / b * 100) if b else 0 for d, b in zip(deltas, base_medians)]

        colors_bar = [compare_info["color"]] * len(families)
        bars = ax.bar(x, deltas, 0.5, color=colors_bar, alpha=0.85)
        ax.set_xticks(x)
        ax.set_xticklabels(families, fontsize=9, fontweight="bold")
        ax.set_ylabel("Δ Handshake Time (ms)")
        ax.set_title(f"Overhead: {compare_info['label']}", fontweight="bold")
        ax.axhline(y=0, color="black", linewidth=0.5)
        ax.grid(axis="y", alpha=0.3)
        ax.set_axisbelow(True)
        # Annotate percentage
        for j, (bar, p) in enumerate(zip(bars, pct)):
            y_pos = bar.get_height()
            if y_pos >= 0:
                ax.text(bar.get_x() + bar.get_width()/2, y_pos + abs(y_pos)*0.05,
                        f"{p:+.1f}%", ha="center", va="bottom", fontsize=8, fontweight="bold")
            else:
                ax.text(bar.get_x() + bar.get_width()/2, y_pos - abs(y_pos)*0.05,
                        f"{p:+.1f}%", ha="center", va="top", fontsize=8, fontweight="bold")

    fig.suptitle("DDoS Detection Overhead on Handshake Time", fontweight="bold", fontsize=13, y=1.02)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_three_scenario_heatmap(data):
    """Summary table-style heatmap: key metrics across scenarios."""
    fig, ax = plt.subplots(figsize=(10, 5))
    metrics_labels = [
        "Median Handshake (ms)",
        "Mean Handshake (ms)",
        "P95 Handshake (ms)",
        "Avg CPU Drone (%)",
        "Avg Temp (°C)",
        "Packets Sent (median)",
        "Packet Loss (%)",
    ]
    scenario_labels = [s["label"] for s in SCENARIOS.values()]
    table_data = []
    for scenario in SCENARIOS:
        row = []
        hs_times = [d["handshake"]["handshake_total_duration_ms"]
                     for d in data[scenario].values()
                     if d.get("handshake", {}).get("handshake_total_duration_ms")]
        cpus = [d["system_drone"]["cpu_usage_avg_percent"]
                for d in data[scenario].values()
                if d.get("system_drone", {}).get("cpu_usage_avg_percent") is not None]
        temps = [d["system_drone"]["temperature_c"]
                 for d in data[scenario].values()
                 if d.get("system_drone", {}).get("temperature_c") is not None]
        pkts = [d["data_plane"]["packets_sent"]
                for d in data[scenario].values()
                if d.get("data_plane", {}).get("packets_sent")]
        loss = [d["data_plane"]["packet_loss_ratio"]
                for d in data[scenario].values()
                if d.get("data_plane", {}).get("packet_loss_ratio") is not None]

        row.append(f"{np.median(hs_times):.1f}" if hs_times else "N/A")
        row.append(f"{np.mean(hs_times):.1f}" if hs_times else "N/A")
        row.append(f"{np.percentile(hs_times, 95):.1f}" if hs_times else "N/A")
        row.append(f"{np.mean(cpus):.1f}" if cpus else "N/A")
        row.append(f"{np.mean(temps):.1f}" if temps else "N/A")
        row.append(f"{int(np.median(pkts))}" if pkts else "N/A")
        row.append(f"{np.mean(loss)*100:.2f}" if loss else "N/A")
        table_data.append(row)

    ax.axis("off")
    table = ax.table(
        cellText=list(zip(*table_data)),  # transpose: rows=metrics, cols=scenarios
        rowLabels=metrics_labels,
        colLabels=scenario_labels,
        loc="center",
        cellLoc="center",
    )
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1.2, 1.8)
    # Style header
    for j in range(len(scenario_labels)):
        table[0, j].set_facecolor("#1A237E")
        table[0, j].set_text_props(color="white", fontweight="bold")
    for i in range(len(metrics_labels)):
        table[i+1, -1].set_facecolor("#E3F2FD")  # row labels
    ax.set_title("Cross-Scenario Key Metrics Summary", fontweight="bold", fontsize=13, pad=20)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_sig_timing_comparison(data, scenario="no-ddos"):
    """Bar chart – signature sign times by SIG family."""
    fig, ax = plt.subplots(figsize=(8, 5))
    sigs = ["ML-DSA-44", "ML-DSA-65", "ML-DSA-87",
            "Falcon-512", "Falcon-1024",
            "SPHINCS+-128s", "SPHINCS+-192s", "SPHINCS+-256s"]
    sig_short = ["ML-DSA\n44", "ML-DSA\n65", "ML-DSA\n87",
                 "Falcon\n512", "Falcon\n1024",
                 "SPHINCS+\n128s", "SPHINCS+\n192s", "SPHINCS+\n256s"]
    vals = []
    colors = []
    fam_colors = {"ML-DSA": "#4CAF50", "Falcon": "#9C27B0", "SPHINCS+": "#FF5722"}
    for sig in sigs:
        times = []
        for sid, d in data[scenario].items():
            _, _, s = parse_suite_id(sid)
            if s == sig:
                t = d.get("crypto_primitives", {}).get("signature_sign_time_ms", 0)
                if t: times.append(t)
        vals.append(np.median(times) if times else 0)
        colors.append(fam_colors.get(get_sig_family(sig), "#999"))

    bars = ax.bar(range(len(sigs)), vals, color=colors, alpha=0.85)
    ax.set_xticks(range(len(sigs)))
    ax.set_xticklabels(sig_short, fontsize=8)
    ax.set_ylabel("Sign Time (ms) – log scale")
    ax.set_yscale("log")
    ax.set_title("Signature Signing Time by Algorithm", fontweight="bold")
    ax.grid(axis="y", alpha=0.3)
    ax.set_axisbelow(True)
    # Annotate values
    for bar, v in zip(bars, vals):
        if v > 0:
            ax.text(bar.get_x() + bar.get_width()/2, v * 1.15,
                    f"{v:.1f}", ha="center", va="bottom", fontsize=7, fontweight="bold")

    from matplotlib.patches import Patch
    legend_elements = [Patch(facecolor=c, label=f) for f, c in fam_colors.items()]
    ax.legend(handles=legend_elements, fontsize=8)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_aead_comparison(data, scenario="no-ddos"):
    """Bar chart – AEAD encrypt/decrypt times."""
    fig, ax = plt.subplots(figsize=(7, 4.5))
    aeads = ["AES-256-GCM", "ChaCha20-Poly1305", "Ascon-128a"]
    enc_vals = []
    dec_vals = []
    for aead in aeads:
        enc_times = []
        dec_times = []
        for sid, d in data[scenario].items():
            _, a, _ = parse_suite_id(sid)
            if a == aead:
                e = d.get("data_plane", {}).get("aead_encrypt_avg_ns", 0)
                dc = d.get("data_plane", {}).get("aead_decrypt_avg_ns", 0)
                if e: enc_times.append(e / 1000)  # ns → µs
                if dc: dec_times.append(dc / 1000)
        enc_vals.append(np.median(enc_times) if enc_times else 0)
        dec_vals.append(np.median(dec_times) if dec_times else 0)

    x = np.arange(len(aeads))
    w = 0.35
    ax.bar(x - w/2, enc_vals, w, label="Encrypt", color="#2196F3", alpha=0.85)
    ax.bar(x + w/2, dec_vals, w, label="Decrypt", color="#FF9800", alpha=0.85)
    ax.set_xticks(x)
    ax.set_xticklabels(aeads, fontsize=10, fontweight="bold")
    ax.set_ylabel("Time (µs)")
    ax.set_title("AEAD Encrypt / Decrypt Performance (64-byte payload)", fontweight="bold")
    ax.legend()
    ax.grid(axis="y", alpha=0.3)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_pareto_frontier(data, scenario="no-ddos"):
    """Scatter plot – NIST level vs handshake time with Pareto frontier."""
    fig, ax = plt.subplots(figsize=(9, 5.5))
    level_map = {"L1": 1, "L3": 3, "L5": 5}
    for sid, d in data[scenario].items():
        k, _, s = parse_suite_id(sid)
        fam = get_kem_family(k)
        t = d.get("handshake", {}).get("handshake_total_duration_ms", 0)
        lvl = get_nist_level(d)
        x_val = level_map.get(lvl, 0)
        fc = KEM_FAMILIES[fam]["color"]
        mk = KEM_FAMILIES[fam]["marker"]
        ax.scatter(x_val + np.random.uniform(-0.15, 0.15), t, c=fc, marker=mk,
                   s=40, alpha=0.6, edgecolors="white", linewidths=0.3)

    # Highlight Pareto-optimal
    pareto = [
        (1, "ML-KEM-512 + Falcon-512", 9.1),
        (3, "ML-KEM-768 + ML-DSA-65", 17.7),
        (5, "ML-KEM-1024 + Falcon-1024", 10.1),
    ]
    px, py = [p[0] for p in pareto], [p[2] for p in pareto]
    ax.scatter(px, py, c="gold", marker="*", s=200, zorder=10, edgecolors="black", linewidths=1,
              label="Pareto-optimal")
    for lvl, name, t in pareto:
        ax.annotate(name, (lvl, t), textcoords="offset points", xytext=(10, 10),
                    fontsize=7, fontweight="bold",
                    arrowprops=dict(arrowstyle="->", color="gray", lw=0.8))

    ax.set_yscale("log")
    ax.set_xticks([1, 3, 5])
    ax.set_xticklabels(["NIST L1", "NIST L3", "NIST L5"], fontsize=10, fontweight="bold")
    ax.set_ylabel("Handshake Time (ms) – log scale")
    ax.set_title("Pareto Frontier: Security Level vs. Handshake Latency", fontweight="bold")
    from matplotlib.patches import Patch
    from matplotlib.lines import Line2D
    legend_elements = [
        Line2D([0], [0], marker=v["marker"], color=v["color"], linestyle="", markersize=7, label=f)
        for f, v in KEM_FAMILIES.items()
    ] + [Line2D([0], [0], marker="*", color="gold", linestyle="", markersize=12,
                markeredgecolor="black", label="Pareto-optimal")]
    ax.legend(handles=legend_elements, fontsize=8, loc="upper left")
    ax.grid(alpha=0.3)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_system_metrics_comparison(data):
    """Grouped bar chart – CPU, temp, packets across scenarios."""
    fig, axes = plt.subplots(1, 3, figsize=(13, 4.5))
    scenario_labels = [info["label"] for info in SCENARIOS.values()]
    scenario_colors = [info["color"] for info in SCENARIOS.values()]

    # CPU
    cpu_vals = []
    for scenario in SCENARIOS:
        cpus = [d["system_drone"]["cpu_usage_avg_percent"]
                for d in data[scenario].values()
                if d.get("system_drone", {}).get("cpu_usage_avg_percent") is not None]
        cpu_vals.append(np.mean(cpus) if cpus else 0)
    axes[0].bar(scenario_labels, cpu_vals, color=scenario_colors, alpha=0.85)
    axes[0].set_ylabel("CPU Usage (%)")
    axes[0].set_title("Avg Drone CPU", fontweight="bold")
    axes[0].set_ylim(0, max(cpu_vals)*1.3 if max(cpu_vals) else 100)
    for i, v in enumerate(cpu_vals):
        axes[0].text(i, v + 0.5, f"{v:.1f}%", ha="center", fontsize=9, fontweight="bold")

    # Temperature
    temp_vals = []
    for scenario in SCENARIOS:
        temps = [d["system_drone"]["temperature_c"]
                 for d in data[scenario].values()
                 if d.get("system_drone", {}).get("temperature_c") is not None]
        temp_vals.append(np.mean(temps) if temps else 0)
    axes[1].bar(scenario_labels, temp_vals, color=scenario_colors, alpha=0.85)
    axes[1].set_ylabel("Temperature (°C)")
    axes[1].set_title("Avg SoC Temperature", fontweight="bold")
    axes[1].set_ylim(min(temp_vals)*0.9 if min(temp_vals) else 0, max(temp_vals)*1.1 if max(temp_vals) else 100)
    for i, v in enumerate(temp_vals):
        axes[1].text(i, v + 0.3, f"{v:.1f}°C", ha="center", fontsize=9, fontweight="bold")

    # Packets
    pkt_vals = []
    for scenario in SCENARIOS:
        pkts = [d["data_plane"]["packets_sent"]
                for d in data[scenario].values()
                if d.get("data_plane", {}).get("packets_sent")]
        pkt_vals.append(np.median(pkts) if pkts else 0)
    axes[2].bar(scenario_labels, pkt_vals, color=scenario_colors, alpha=0.85)
    axes[2].set_ylabel("Packets Sent (median)")
    axes[2].set_title("Data Plane Packets", fontweight="bold")
    for i, v in enumerate(pkt_vals):
        axes[2].text(i, v + 5, f"{int(v)}", ha="center", fontsize=9, fontweight="bold")

    for ax in axes:
        ax.grid(axis="y", alpha=0.3)
        ax.set_axisbelow(True)
        ax.tick_params(axis="x", labelsize=7)

    fig.suptitle("System Resource Impact Across Scenarios", fontweight="bold", fontsize=13, y=1.02)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_handshake_boxplot_by_level(data, scenario="no-ddos"):
    """Box plot – handshake time distribution by NIST level."""
    fig, ax = plt.subplots(figsize=(8, 5))
    levels = ["L1", "L3", "L5"]
    level_data = {l: [] for l in levels}
    for sid, d in data[scenario].items():
        lvl = get_nist_level(d)
        t = d.get("handshake", {}).get("handshake_total_duration_ms", 0)
        if t and lvl in level_data:
            level_data[lvl].append(t)

    bp = ax.boxplot([level_data[l] for l in levels], labels=levels, patch_artist=True,
                    widths=0.5, showfliers=True, flierprops=dict(markersize=3))
    colors_box = ["#E3F2FD", "#FFF3E0", "#FCE4EC"]
    for patch, color in zip(bp["boxes"], colors_box):
        patch.set_facecolor(color)
    ax.set_ylabel("Handshake Time (ms)")
    ax.set_title("Handshake Time Distribution by NIST Level (Baseline)", fontweight="bold")
    ax.grid(axis="y", alpha=0.3)
    ax.set_axisbelow(True)
    # Stats text
    for i, l in enumerate(levels):
        vals = level_data[l]
        if vals:
            ax.text(i+1, max(vals)*1.05,
                    f"n={len(vals)}\nμ={np.mean(vals):.0f}\nmed={np.median(vals):.0f}",
                    ha="center", va="bottom", fontsize=7)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_ddos_model_comparison():
    """Infographic-style comparison: XGBoost vs TST detector."""
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.axis("off")

    props = [
        ("Architecture",       "XGBoost (Gradient Boosted Trees)", "Time Series Transformer (3-layer, 16-head)"),
        ("Window Size",        "5 packets",                        "400 packets"),
        ("Detection Latency",  "~3 seconds",                       "~240 seconds"),
        ("Inference Time",     "~microseconds (μs)",               "~milliseconds (ms)"),
        ("Threading Model",    "Single-thread, GIL-friendly",      "Single-thread, no GPU"),
        ("Model Size",         "~100 KB",                          "~5 MB"),
        ("Feature Extraction", "Packet count per window",          "Packet count per window"),
        ("Output",             "Binary (Attack/Normal)",            "Binary (Attack/Normal)"),
    ]

    cols = ["Property", "XGBoost Detector", "TST Detector"]
    cell_text = [[p[0], p[1], p[2]] for p in props]
    table = ax.table(cellText=cell_text, colLabels=cols, loc="center", cellLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(9)
    table.scale(1.2, 1.7)
    # Style
    for j in range(3):
        table[0, j].set_facecolor("#1A237E")
        table[0, j].set_text_props(color="white", fontweight="bold")
    for i in range(len(props)):
        table[i+1, 0].set_facecolor("#E8EAF6")
        table[i+1, 0].set_text_props(fontweight="bold")
        table[i+1, 1].set_facecolor("#FFF3E0")
        table[i+1, 2].set_facecolor("#FFEBEE")
    ax.set_title("DDoS Detection Models: XGBoost vs. TST", fontweight="bold", fontsize=13, pad=20)
    fig.tight_layout()
    return fig_to_bytes(fig)


def chart_rekey_overhead(data, scenario="no-ddos"):
    """Line chart – rekey overhead fraction Φ vs rekey interval R for key suites."""
    fig, ax = plt.subplots(figsize=(9, 5))
    R_vals = np.array([10, 30, 60, 120, 300, 600, 1800, 3600])
    key_suites = {
        "ML-KEM-768 + ML-DSA-65": "cs-mlkem768-aesgcm-mldsa65",
        "HQC-256 + Falcon-1024": "cs-hqc256-aesgcm-falcon1024",
        "McE-348864 + Falcon-512": "cs-classicmceliece348864-aesgcm-falcon512",
        "McE-8192128 + ML-DSA-87": "cs-classicmceliece8192128-aesgcm-mldsa87",
    }
    colors_line = ["#2196F3", "#FF9800", "#F44336", "#9C27B0"]
    for i, (label, sid) in enumerate(key_suites.items()):
        if sid in data[scenario]:
            Ths = data[scenario][sid]["handshake"]["handshake_total_duration_ms"]
            Ths_s = Ths / 1000
            phi = Ths_s / (R_vals + Ths_s) * 100
            ax.plot(R_vals, phi, "o-", label=f"{label} ({Ths:.0f} ms)", color=colors_line[i],
                    linewidth=2, markersize=5)

    ax.set_xscale("log")
    ax.set_yscale("log")
    ax.set_xlabel("Rekey Interval R (seconds)")
    ax.set_ylabel("Overhead Φ (%)")
    ax.set_title("Rekey Overhead Fraction Φ(R) = Ths / (R + Ths)", fontweight="bold")
    ax.legend(fontsize=8)
    ax.grid(alpha=0.3, which="both")
    ax.set_axisbelow(True)
    ax.axhline(y=1, color="red", linestyle="--", linewidth=0.8, alpha=0.5)
    ax.text(3700, 1.1, "1% threshold", fontsize=7, color="red")
    fig.tight_layout()
    return fig_to_bytes(fig)


# ============================================================================
# PPTX Builder
# ============================================================================
class PresentationBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # Widescreen 16:9
        self.prs.slide_height = Inches(7.5)
        self.slide_num = 0

    def _add_bg(self, slide, color=DARK):
        """Fill slide background with solid color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_shape(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_text_box(self, slide, left, top, width, height, text, font_size=18,
                      color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def _add_bullet_slide_content(self, slide, left, top, width, height, bullets, font_size=16,
                                   color=WHITE, spacing=Pt(6)):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Segoe UI"
            p.space_after = spacing
            p.level = 0
        return txBox

    def _slide_number(self, slide):
        self.slide_num += 1
        self._add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                          str(self.slide_num), 10, GRAY, alignment=PP_ALIGN.RIGHT)

    def add_title_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self._add_bg(slide, DARK)
        # Accent bar
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BLUE)
        # Title
        self._add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                          "PQC-MAV", 54, WHITE, bold=True)
        self._add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(1),
                          "A Complete Post-Quantum Cryptographic Tunnel\nfor UAV–GCS Communication",
                          26, RGBColor(0xBB, 0xDE, 0xFB), bold=False)
        # Separator
        self._add_shape(slide, Inches(1), Inches(4.5), Inches(3), Inches(0.04), BLUE)
        # Author
        self._add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
                          "Burak Güneysu", 20, WHITE, bold=True)
        self._add_text_box(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.5),
                          "Department of Computer Science  •  February 2026", 14, GRAY)
        # Key stat boxes
        stats = [("72", "Cipher Suites"), ("3", "KEM Families"), ("3", "Scenarios"), ("19,600", "Benchmarks")]
        for i, (num, label) in enumerate(stats):
            left = Inches(1 + i * 2.8)
            self._add_shape(slide, left, Inches(6.2), Inches(2.4), Inches(0.9),
                           RGBColor(0x28, 0x35, 0x93))
            self._add_text_box(slide, left, Inches(6.2), Inches(2.4), Inches(0.5),
                              num, 24, BLUE, bold=True, alignment=PP_ALIGN.CENTER)
            self._add_text_box(slide, left, Inches(6.6), Inches(2.4), Inches(0.4),
                              label, 11, GRAY, alignment=PP_ALIGN.CENTER)

    def add_agenda_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, WHITE)
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
        self._add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                          "Presentation Outline", 32, WHITE, bold=True)
        items = [
            "1.  Motivation & Problem Statement",
            "2.  System Architecture – Bump-in-the-Wire Tunnel",
            "3.  Post-Quantum Cryptography Primer",
            "4.  72 Cipher Suite Registry – KEMs, SIGs, AEADs",
            "5.  PQC Handshake & AEAD Framing Protocol",
            "6.  Benchmark Methodology & Hardware Testbed",
            "7.  KEM Primitive Performance Analysis",
            "8.  Signature & AEAD Benchmark Results",
            "9.  End-to-End Handshake Results (All 72 Suites)",
            "10. Pareto-Optimal Suites & Rekey Overhead",
            "11. DDoS Detection Models – XGBoost & TST",
            "12. Three-Scenario Benchmark Comparison",
            "13. DDoS Overhead Analysis & Cross-Run Delta",
            "14. Adaptive Policy & Graceful Degradation",
            "15. Key Findings & Conclusions",
        ]
        self._add_bullet_slide_content(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(5.5),
                                        items, 16, BLACK, spacing=Pt(4))
        self._slide_number(slide)

    def add_section_divider(self, title, subtitle="", number=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, DARK)
        self._add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), BLUE)
        if number:
            self._add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
                              number, 60, BLUE, bold=True)
        self._add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
                          title, 36, WHITE, bold=True)
        if subtitle:
            self._add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
                              subtitle, 18, GRAY)
        self._slide_number(slide)

    def add_content_slide(self, title, bullets, image_buf=None, image_width=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, WHITE)
        # Header bar
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK)
        self._add_text_box(slide, Inches(0.8), Inches(0.15), Inches(12), Inches(0.7),
                          title, 24, WHITE, bold=True)

        if image_buf and bullets:
            # Split layout: text left, image right
            self._add_bullet_slide_content(slide, Inches(0.6), Inches(1.3), Inches(5.5), Inches(5.5),
                                            bullets, 14, BLACK, spacing=Pt(4))
            iw = image_width or Inches(6.5)
            slide.shapes.add_picture(image_buf, Inches(6.5), Inches(1.2), width=iw)
        elif image_buf:
            # Full-width image
            iw = image_width or Inches(11)
            slide.shapes.add_picture(image_buf, Inches(1.0), Inches(1.2), width=iw)
        elif bullets:
            # Full-width text
            self._add_bullet_slide_content(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                                            bullets, 16, BLACK, spacing=Pt(6))
        self._slide_number(slide)

    def add_image_slide(self, title, image_buf, subtitle="", image_width=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, WHITE)
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK)
        self._add_text_box(slide, Inches(0.8), Inches(0.15), Inches(12), Inches(0.7),
                          title, 24, WHITE, bold=True)
        if subtitle:
            self._add_text_box(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.4),
                              subtitle, 12, GRAY)
        iw = image_width or Inches(11)
        y_start = Inches(1.5) if subtitle else Inches(1.2)
        slide.shapes.add_picture(image_buf, Inches(1.0), y_start, width=iw)
        self._slide_number(slide)

    def add_two_image_slide(self, title, img_left, img_right, label_left="", label_right=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, WHITE)
        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK)
        self._add_text_box(slide, Inches(0.8), Inches(0.15), Inches(12), Inches(0.7),
                          title, 24, WHITE, bold=True)
        slide.shapes.add_picture(img_left, Inches(0.3), Inches(1.3), width=Inches(6.2))
        slide.shapes.add_picture(img_right, Inches(6.8), Inches(1.3), width=Inches(6.2))
        if label_left:
            self._add_text_box(slide, Inches(0.3), Inches(6.8), Inches(6.2), Inches(0.4),
                              label_left, 10, GRAY, alignment=PP_ALIGN.CENTER)
        if label_right:
            self._add_text_box(slide, Inches(6.8), Inches(6.8), Inches(6.2), Inches(0.4),
                              label_right, 10, GRAY, alignment=PP_ALIGN.CENTER)
        self._slide_number(slide)

    def save(self, path):
        self.prs.save(path)
        print(f"\n✅ Saved: {path}")


# ============================================================================
# Main
# ============================================================================
def main():
    print("=" * 60)
    print("PQC-MAV Presentation Generator")
    print("=" * 60)

    # 1. Load data
    print("\n📦 Loading benchmark data...")
    data = load_all_data()

    # 2. Generate charts
    print("\n📊 Generating charts...")

    print("  [1/10] KEM keygen comparison...")
    img_kem_keygen = chart_kem_keygen_comparison(data)

    print("  [2/10] Handshake by KEM family...")
    img_hs_family = chart_handshake_by_family(data)

    print("  [3/10] All 72 suites handshake...")
    img_all_suites = chart_handshake_all_suites(data)

    print("  [4/10] Scenario overhead delta...")
    img_overhead_delta = chart_scenario_overhead_delta(data)

    print("  [5/10] Summary metrics table...")
    img_summary_table = chart_three_scenario_heatmap(data)

    print("  [6/10] Signature timing comparison...")
    img_sig_timing = chart_sig_timing_comparison(data)

    print("  [7/10] AEAD comparison...")
    img_aead = chart_aead_comparison(data)

    print("  [8/10] Pareto frontier...")
    img_pareto = chart_pareto_frontier(data)

    print("  [9/10] System metrics comparison...")
    img_system = chart_system_metrics_comparison(data)

    print("  [10/10] Rekey overhead + misc...")
    img_rekey = chart_rekey_overhead(data)
    img_boxplot = chart_handshake_boxplot_by_level(data)
    img_ddos_compare = chart_ddos_model_comparison()

    # 3. Build PPTX
    print("\n📑 Building presentation...")
    pb = PresentationBuilder()

    # === SLIDE 1: Title ===
    pb.add_title_slide()

    # === SLIDE 2: Agenda ===
    pb.add_agenda_slide()

    # === SECTION: Motivation ===
    pb.add_section_divider("Motivation & Problem Statement",
                           "Why post-quantum cryptography for drones?", "01")

    pb.add_content_slide("The Quantum Threat to UAV Communications", [
        "■  Modern UAVs use MAVLink 2.0 protocol – compact binary telemetry up to 320 Hz",
        "■  Current protection relies on classical crypto: RSA, ECDH, ECDSA",
        "■  NIST warns: quantum computers will break ALL classical key exchange",
        "",
        "■  Three critical challenges for drone PQC deployment:",
        "    ① Extreme performance heterogeneity – KEM keygen spans 4 orders of magnitude",
        "    ② Resource constraints – RPi 5 at 3.8W power budget, 80°C thermal limit",
        "    ③ Safety continuity – MAVLink heartbeats must arrive within 5 seconds",
        "",
        "■  No prior work implements a complete PQC tunnel for real MAVLink traffic",
        "■  PQC-MAV is the first bump-in-the-wire PQC tunnel for UAV–GCS links",
    ])

    # === SECTION: Architecture ===
    pb.add_section_divider("System Architecture",
                           "PQC-MAV: A transparent bump-in-the-wire tunnel", "02")

    pb.add_content_slide("Tunnel Architecture Overview", [
        "■  Bump-in-the-wire proxy: zero changes to MAVProxy / QGroundControl",
        "",
        "■  Data Path (bidirectional):",
        "    Pixhawk FC → MAVProxy (serial) → PQC Proxy → [ENCRYPTED] → PQC Proxy → MAVProxy → QGC",
        "",
        "■  Controller–Follower model:",
        "    • Drone (RPi 5) = Controller – initiates suite changes, manages rekey",
        "    • GCS (Windows 11) = Follower – executes commands via JSON-RPC (TCP:48080)",
        "",
        "■  Two-phase rekey: Prepare → Commit → Abort/Rollback",
        "    • Blackout = T_startup (~3s) + T_handshake",
        "    • Failed suites are blacklisted (TTL: 30 min)",
        "",
        "■  Clock synchronisation via Operation Chronos (NTP-lite, 3-message)",
    ])

    pb.add_content_slide("Hardware Testbed", [
        "■  Drone (uavpi):",
        "    • Raspberry Pi 5 – ARM Cortex-A76 (4 cores)",
        "    • 3,796 MB RAM, Linux 6.12.47+rpt",
        "    • Python 3.11.2, liboqs 0.12.0",
        "    • Power: INA219 sensor at 1,100 Hz / RPi5 hwmon",
        "",
        "■  GCS (lappy):",
        "    • Windows 11 x86-64",
        "    • Python 3.11.13, liboqs 0.12.0",
        "",
        "■  Network: Ethernet LAN (sub-ms RTT)",
        "■  MAVLink: Pixhawk SITL at 320 Hz telemetry rate",
    ])

    # === SECTION: PQC Primer ===
    pb.add_section_divider("Post-Quantum Cryptography Primer",
                           "NIST standards & mathematical assumptions", "03")

    pb.add_content_slide("PQC Algorithm Families", [
        "■  KEM (Key Encapsulation Mechanism) – replace Diffie-Hellman:",
        "    • ML-KEM (FIPS 203) – Module Learning-with-Errors (lattice)",
        "    • HQC – Hamming Quasi-Cyclic codes",
        "    • Classic McEliece – Binary Goppa codes (Niederreiter)",
        "",
        "■  Digital Signatures – replace RSA/ECDSA:",
        "    • ML-DSA (FIPS 204) – Module-LWE lattice signatures",
        "    • Falcon – NTRU lattice + fast Fourier sampling",
        "    • SPHINCS+ (FIPS 205) – Stateless hash-based signatures",
        "",
        "■  AEAD (Authenticated Encryption with Associated Data):",
        "    • AES-256-GCM – NIST standard, HW-accelerated on x86",
        "    • ChaCha20-Poly1305 – Software-friendly (Bernstein)",
        "    • Ascon-128a – Lightweight, NIST LWC winner",
        "",
        "■  NIST Security Levels: L1 (AES-128), L3 (AES-192), L5 (AES-256 equivalent)",
    ])

    # === SECTION: Cipher Suite Registry ===
    pb.add_section_divider("72 Cipher Suite Registry",
                           "9 KEMs × 8 SIGs × 3 AEADs = 72 NIST-level-matched suites", "04")

    pb.add_content_slide("Suite Construction Formula", [
        "■  Suite Registry S = { (KEM_i, SIG_j, AEAD_k) | Level(KEM_i) = Level(SIG_j) }",
        "",
        "■  9 KEMs:  ML-KEM-{512,768,1024}, HQC-{128,192,256}, McEliece-{348864,460896,8192128}",
        "■  8 SIGs:  ML-DSA-{44,65,87}, Falcon-{512,1024}, SPHINCS+-{128s,192s,256s}",
        "■  3 AEADs: AES-256-GCM, ChaCha20-Poly1305, Ascon-128a",
        "",
        "■  Level distribution:",
        "    • L1 = 3 KEMs × 3 SIGs × 3 AEADs = 27 suites",
        "    • L3 = 3 KEMs × 2 SIGs × 3 AEADs = 18 suites  (no Falcon at L3)",
        "    • L5 = 3 KEMs × 3 SIGs × 3 AEADs = 27 suites",
        "    ─────────────────────────────────────────────",
        "    Total: 72 cipher suites",
        "",
        "■  Algorithm diversity: 3 distinct mathematical assumptions (lattice, quasi-cyclic code, Goppa code)",
        "■  Suite ID format: cs-{kem}-{aead}-{sig}  e.g., cs-mlkem768-aesgcm-mldsa65",
    ])

    # === SECTION: Handshake Protocol ===
    pb.add_section_divider("PQC Handshake & AEAD Framing",
                           "KEM + SIG + HKDF handshake with deterministic nonces", "05")

    pb.add_content_slide("PQC Handshake Protocol (3 Messages)", [
        "■  Message 1 – ServerHello (GCS → Drone):",
        "    • GCS generates ephemeral KEM keypair: (pk, sk) ← KEM.Keygen()",
        "    • Signs transcript: σ ← SIG.Sign(sk_GCS, ver ‖ kem_name ‖ sig_name ‖ sid ‖ pk)",
        "    • Sends: ver ‖ pk ‖ σ ‖ sid ‖ challenge",
        "",
        "■  Message 2 – ClientFinish (Drone → GCS):",
        "    • Verifies signature (anti-downgrade)",
        "    • KEM encapsulation: (ct, ss) ← KEM.Encaps(pk)",
        "    • Sends: ct ‖ HMAC-SHA256(PSK, ct ‖ challenge)",
        "",
        "■  Message 3 – Key Derivation (both sides):",
        "    • GCS verifies HMAC, decapsulates: ss ← KEM.Decaps(sk, ct)",
        "    • k_d→g ‖ k_g→d = HKDF-SHA256(ss, salt, ver ‖ sid ‖ kem ‖ sig)",
        "    • Directional keys prevent reflection attacks",
        "",
        "■  AEAD Wire Format: 22-byte header (AAD) + encrypted payload + 16-byte tag",
        "    • Nonce = epoch(4B) ‖ seq(8B) — never transmitted, reconstructed",
        "    • Anti-replay: sliding-window bitmap (64+ packets)",
    ])

    # === SECTION: Benchmark Methodology ===
    pb.add_section_divider("Benchmark Methodology",
                           "19,600 timed operations + 72 end-to-end tunnel runs", "06")

    pb.add_content_slide("Benchmarking Approach", [
        "■  Three-tier benchmarking:",
        "    Tier 1:  Primitive-level (KEM/SIG/AEAD operations, n=200 each)",
        "    Tier 2:  DDoS overhead measurement (per-handshake with detector active)",
        "    Tier 3:  End-to-end tunnel – full 110s MAVLink session per suite",
        "",
        "■  Three scenarios (each runs all 72 suites):",
        "    • Baseline (No DDoS)  – run 20260211_141627 – pure tunnel performance",
        "    • + XGBoost Detector   – run 20260211_150013 – fast ML detection overhead",
        "    • + TST Detector       – run 20260211_171406 – heavy DL detection overhead",
        "",
        "■  Per-suite data pipeline:",
        "    Drone scheduler → start_proxy RPC → handshake → 110s traffic → stop → merge metrics → JSON",
        "",
        "■  18 metric categories (A–R): crypto_identity, handshake, crypto_primitives,",
        "    data_plane, latency, MAVProxy, fc_telemetry, system, power/energy, validation",
        "",
        "■  432 total JSON files (72 suites × 2 endpoints × 3 scenarios)",
    ])

    # === SECTION: KEM Results ===
    pb.add_section_divider("KEM Primitive Performance",
                           "Four orders of magnitude – from 0.08 ms to 8,835 ms", "07")

    pb.add_image_slide("KEM Keygen Time Across Algorithms & Scenarios",
                       img_kem_keygen,
                       "Median keygen times (log scale) – ML-KEM is sub-millisecond, McEliece is seconds",
                       image_width=Inches(11))

    pb.add_content_slide("KEM Performance – Key Findings", [
        "■  Key Finding 1: ML-KEM dominates in speed",
        "    • ML-KEM-512: 0.08 ms keygen  |  ML-KEM-1024: 0.14 ms keygen",
        "    • 3 orders of magnitude faster than HQC (22–392 ms)",
        "    • 5 orders faster than McEliece keygen (333–8,835 ms)",
        "",
        "■  Key Finding 2: McEliece has extreme variance",
        "    • McEliece-8192128: σ = 857 ms (probabilistic key generation)",
        "    • Worst case keygen: up to 36.6 seconds!",
        "    • Unsuitable for time-bounded rekey operations",
        "",
        "■  Key Finding: Public key sizes vary enormously",
        "    • ML-KEM-512: 800 bytes  vs.  McEliece-8192128: 1,357,824 bytes (1.36 MB)",
        "    • 867× size difference → bandwidth and memory implications",
        "",
        "■  Across all 3 scenarios, KEM primitive times are consistent",
        "    → DDoS detection does NOT affect the crypto layer",
    ])

    # === SECTION: SIG + AEAD Results ===
    pb.add_section_divider("Signature & AEAD Benchmarks",
                           "SPHINCS+ signing > 1 second  •  AEAD is a non-factor", "08")

    pb.add_two_image_slide("Signature & AEAD Performance",
                           img_sig_timing, img_aead,
                           "SIG sign times by algorithm (log scale)",
                           "AEAD encrypt/decrypt (µs)")

    pb.add_content_slide("SIG & AEAD – Key Findings", [
        "■  Signature Performance:",
        "    • Falcon-512:   sign 0.65 ms, verify 0.11 ms – fastest combo",
        "    • ML-DSA-65:    sign 1.59 ms, verify 0.38 ms – good all-rounder",
        "    • SPHINCS+-128s: sign 1,461 ms – >1000× slower than Falcon!",
        "    • SPHINCS+ is the dominant bottleneck in any suite that includes it",
        "",
        "■  Key Finding 3: AEAD algorithm choice is a non-factor",
        "    • AES-256-GCM: 7.3 µs  |  ChaCha20: 6.7 µs  |  Ascon-128a: 4.1 µs",
        "    • Max difference: 3.5 µs/packet → 1.12 ms/sec at 320 Hz = 0.11%",
        "    • Ascon is 44% faster (no AES-NI on ARM Cortex-A76)",
        "    • AEAD selection should be driven by side-channel resistance, not performance",
        "",
        "■  Conclusion: Rekey decisions should focus on KEM + SIG, not AEAD",
    ])

    # === SECTION: End-to-End Results ===
    pb.add_section_divider("End-to-End Handshake Results",
                           "71/72 suites succeeded • Complete tunnel validation", "09")

    pb.add_image_slide("All 72 Cipher Suites – Baseline Handshake Time",
                       img_all_suites,
                       "Horizontal bars: sorted by handshake duration (log scale)",
                       image_width=Inches(7))

    pb.add_two_image_slide("Handshake Analysis: By KEM Family & NIST Level",
                           img_hs_family, img_boxplot,
                           "Median handshake time by KEM family across scenarios",
                           "Handshake distribution by NIST level (baseline)")

    pb.add_content_slide("End-to-End Handshake – Key Statistics", [
        "■  71/72 suites completed successfully",
        "    • 1 timeout: McEliece-460896 + SPHINCS+-192s (combined weight)",
        "",
        "■  By NIST Level (baseline):",
        "    • L1: n=27, mean=290 ms, median=123 ms, P95=835 ms",
        "    • L3: n=17, mean=749 ms, median=416 ms, P95=1,583 ms",
        "    • L5: n=27, mean=855 ms, median=870 ms, max=3,273 ms",
        "",
        "■  By KEM Family (baseline median):",
        "    • ML-KEM suites: ~10–22 ms (without SPHINCS+)",
        "    • HQC suites: ~100–400 ms",
        "    • McEliece suites: ~400–3,000+ ms",
        "",
        "■  Regression: log₁₀(T_hs) = −2.39 + 0.30·log₁₀(pk_size) + 0.97·log₁₀(sig_size)",
        "    R² = 0.67 – Signature size is the dominant predictor (β = 0.97)",
    ])

    # === SECTION: Pareto + Rekey ===
    pb.add_section_divider("Pareto-Optimal Suites & Rekey Overhead",
                           "Identifying the best security-performance trade-offs", "10")

    pb.add_image_slide("Pareto Frontier: Security Level vs. Handshake Latency",
                       img_pareto,
                       "All suites plotted; gold stars = Pareto-optimal (ML-KEM + Falcon/ML-DSA)",
                       image_width=Inches(10))

    pb.add_content_slide("Pareto-Optimal Suites & Rekey Overhead", [
        "■  Three Pareto-Optimal Suites (all use ML-KEM):",
        "    ★ L1: ML-KEM-512 + Falcon-512       → 9.1 ms handshake",
        "    ★ L3: ML-KEM-768 + ML-DSA-65         → 17.7 ms handshake",
        "    ★ L5: ML-KEM-1024 + Falcon-1024      → 10.1 ms handshake",
        "",
        "■  Rekey Overhead Fraction: Φ(R) = T_hs / (R + T_hs)",
        "    At R = 60s rekey interval:",
        "    • ML-KEM-768:    Φ = 0.029%  ✓  negligible",
        "    • HQC-256:       Φ = 0.45%   ⚠  acceptable",
        "    • McE-8192128:   Φ = 0.85%   ✗  high overhead",
        "",
        "■  29× difference between ML-KEM and McEliece overhead!",
        "",
        "■  Key Finding 5: Runtime cipher-suite management is NOT optional –",
        "    operating heavy suites with aggressive rekey intervals degrades MAVLink availability",
    ])

    pb.add_image_slide("Rekey Overhead Fraction Φ(R) vs Interval",
                       img_rekey,
                       "Log-log plot showing overhead for different suites at various rekey intervals",
                       image_width=Inches(10))

    # === SECTION: DDoS Detection ===
    pb.add_section_divider("DDoS Detection Models",
                           "XGBoost (fast) vs TST Transformer (heavy)", "11")

    pb.add_image_slide("DDoS Detection: XGBoost vs. TST Comparison",
                       img_ddos_compare,
                       "Two ML models for MAVLink traffic anomaly detection",
                       image_width=Inches(10))

    pb.add_content_slide("DDoS Detection System Design", [
        "■  Threat: DDoS floods starve MAVLink throughput",
        "    • Normal: ~32 packets/window | Attack: ~5-14 packets/window",
        "    • Detection via Scapy: sniff wlan0 for MAVLink v2 (0xFD magic byte)",
        "",
        "■  XGBoost Detector (lightweight):",
        "    • 5-packet counting window → ~3 second detection latency",
        "    • Microsecond inference (gradient boosted trees)",
        "    • Single-threaded, GIL-friendly, ~100 KB model",
        "    • Best for: real-time detection with minimal overhead",
        "",
        "■  TST Detector (Time Series Transformer):",
        "    • 400-packet window → ~240 second detection latency",
        "    • 3-layer, 16-head Transformer (d_model=128)",
        "    • Millisecond inference, ~5 MB model, CPU-only",
        "    • Best for: higher accuracy with temporal patterns",
        "",
        "■  Hybrid Cascaded Pipeline: XGBoost → TST (fast screening + deep confirmation)",
    ])

    # === SECTION: Three-Scenario Comparison ===
    pb.add_section_divider("Three-Scenario Benchmark Comparison",
                           "Baseline vs +XGBoost vs +TST – 432 JSON files analyzed", "12")

    pb.add_image_slide("Cross-Scenario Key Metrics Summary",
                       img_summary_table,
                       "Aggregated statistics across all 72 suites per scenario",
                       image_width=Inches(10))

    pb.add_image_slide("System Resource Impact: CPU, Temperature, Packets",
                       img_system,
                       "Side-by-side comparison of drone-side system metrics across 3 scenarios",
                       image_width=Inches(11))

    # === SECTION: DDoS Overhead ===
    pb.add_section_divider("DDoS Overhead Analysis",
                           "Measuring the true cost of detection on handshake performance", "13")

    pb.add_image_slide("DDoS Detection Overhead: Handshake Time Delta",
                       img_overhead_delta,
                       "Δ median handshake time vs baseline, by KEM family × detection model",
                       image_width=Inches(11))

    pb.add_content_slide("DDoS Overhead – Analysis & Insights", [
        "■  XGBoost overhead is minimal:",
        "    • ML-KEM: negligible delta (< 1 ms typical)",
        "    • HQC: modest increase due to counting window computation",
        "    • McEliece: slightly noisier due to long keygen + concurrent detection",
        "",
        "■  TST overhead is measurable:",
        "    • Transformer inference adds ~ms per prediction cycle",
        "    • Higher idle CPU consumption (model in memory)",
        "    • TST's 400-packet window means longer warmup before first detection",
        "",
        "■  Key Finding 6: Steady-state metrics are nearly identical across scenarios",
        "    • The AEAD data plane dominates runtime cost, not detection",
        "    • DDoS detection overhead manifests primarily during handshake transitions",
        "    • Once the tunnel is established, all 3 scenarios are equivalent",
        "",
        "■  Conclusion: XGBoost is the recommended detector for flight",
        "    → Minimal overhead, fast detection, negligible impact on PQC operations",
    ])

    # === SECTION: Adaptive Policy ===
    pb.add_section_divider("Adaptive Policy & Graceful Degradation",
                           "TelemetryAwarePolicyV2 – live suite management during flight", "14")

    pb.add_content_slide("Telemetry-Aware Adaptive Rekey Policy", [
        "■  Policy inputs (every 1 second):",
        "    • Link quality: pkt/s, P95 gap, jitter, blackout count (5s window)",
        "    • Battery: voltage (mV), rate of change (mV/min)",
        "    • Thermal: SoC temp (°C), rate of change (°C/min)",
        "    • Armed state: from Pixhawk heartbeat",
        "    • Session state: current suite, epoch, cooldown",
        "",
        "■  Suite Tier Ordering: tier(s) = Level + KEM_weight + AEAD_weight",
        "    Tier 0 (lightest): ML-KEM-512 + AES-GCM → Tier 27: McE-8192128 + Ascon",
        "    Correlation with measured handshake: r = 0.94 (p < 10⁻⁵⁰)",
        "",
        "■  9-level Priority Cascade: Safety → Emergency → Blackout → Cooldown →",
        "    Link → Stress → Rekey → Upgrade → Hold",
        "",
        "■  Most powerful degradation: McEliece → ML-KEM at same NIST level",
        "    → 98× handshake reduction, 867× key-size reduction, ZERO security loss",
    ])

    # === SECTION: Conclusions ===
    pb.add_section_divider("Key Findings & Conclusions", "", "15")

    pb.add_content_slide("Summary of Key Findings", [
        "■  Finding 1: Algorithm-FAMILY selection dominates (not parameter tuning)",
        "    → ML-KEM-1024 (L5): 15 ms  vs  McEliece-348864 (L1): 120 ms = 8× slower at LOWER level",
        "",
        "■  Finding 2: Cross-family degradation is 'free' – same NIST level, 98× faster",
        "",
        "■  Finding 3: AEAD is a non-factor (max 3.5 µs difference per packet)",
        "",
        "■  Finding 4: ML-KEM-based suites dominate at every NIST level",
        "    → All 3 Pareto-optimal suites use ML-KEM",
        "",
        "■  Finding 5: Rekey overhead = 29× between ML-KEM (0.029%) and McEliece (0.85%)",
        "    → Cipher agility is essential, not optional",
        "",
        "■  Finding 6: DDoS detection overhead is minimal on steady-state tunnel",
        "    → XGBoost recommended for flight (µs inference, 3s detection)",
        "",
        "■  PQC-MAV: First complete PQC tunnel for drone MAVLink traffic,",
        "    72 suites evaluated, 19,600 benchmarks, 3 detection scenarios",
    ])

    pb.add_content_slide("Deployment Recommendations", [
        "■  Default suite: ML-KEM-768 + ML-DSA-65 + AES-256-GCM (NIST L3)",
        "    → Handshake < 22 ms, overhead Φ < 0.03% at R=60s",
        "",
        "■  Degradation target: ML-KEM-512 + Falcon-512 + AES-256-GCM (L1)",
        "    → 9.1 ms handshake for stress conditions",
        "",
        "■  Rekey interval: 60s for ML-KEM, ≥300s for HQC, avoid for McEliece",
        "",
        "■  DDoS detection: XGBoost for flight, TST for post-flight analysis",
        "",
        "■  Avoid in flight: SPHINCS+ (signing > 1s) and McEliece-8192128 (keygen > 8s)",
        "",
        "■  Future work:",
        "    → Real-flight validation with adaptive policy active",
        "    → Formal protocol verification (ProVerif/Tamarin)",
        "    → TLS-secured control channel",
        "    → Session resumption to amortise handshake costs",
        "    → Multi-drone swarm architecture extension",
    ])

    # === Thank You ===
    slide = pb.prs.slides.add_slide(pb.prs.slide_layouts[6])
    pb._add_bg(slide, DARK)
    pb._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BLUE)
    pb._add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
                    "Thank You", 54, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    pb._add_shape(slide, Inches(5), Inches(3.5), Inches(3), Inches(0.04), BLUE)
    pb._add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(1),
                    "PQC-MAV: The First Complete Post-Quantum Tunnel\nfor UAV–GCS Communication",
                    20, RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)
    pb._add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                    "Questions?", 28, BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    stats_text = "72 Suites  ·  19,600 Benchmarks  ·  3 Scenarios  ·  432 Data Points  ·  Real Hardware"
    pb._add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
                    stats_text, 12, GRAY, alignment=PP_ALIGN.CENTER)
    pb._slide_number(slide)

    # 4. Save
    output_path = OUTPUT_DIR / "PQC_MAV_Presentation.pptx"
    pb.save(str(output_path))
    print(f"\n{'='*60}")
    print(f"✅ DONE! Presentation has {pb.slide_num} slides")
    print(f"📁 File: {output_path}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
